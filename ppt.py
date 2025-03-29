import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL
from pptx.table import _Cell

# 白と黒を基調としたシンプルなカラーパレット
class ColorPalette:
    BACKGROUND = RGBColor(255, 255, 255)  # ホワイト
    TEXT = RGBColor(0, 0, 0)             # ブラック
    ACCENT = RGBColor(50, 50, 50)         # ダークグレー
    LIGHT_ACCENT = RGBColor(220, 220, 220)  # ライトグレー
    FOOTER_BG = RGBColor(240, 240, 240)   # ライトグレー
    FOOTER_TEXT = RGBColor(100, 100, 100) # ミディアムグレー
    HEADING_BG = RGBColor(0, 0, 0)        # ブラック
    HEADING_TEXT = RGBColor(255, 255, 255) # ホワイト
    TABLE_HEADER_BG = RGBColor(40, 40, 40)  # ダークグレー（テーブルヘッダー用）
    TABLE_HEADER_TEXT = RGBColor(255, 255, 255)  # ホワイト（テーブルヘッダー用）
    TABLE_ACCENT_BG = RGBColor(230, 230, 230)  # 薄いグレー（テーブル強調行用）
    TABLE_BORDER = RGBColor(180, 180, 180)  # ミディアムグレー（テーブル罫線用）

# モダンで洗練されたフォント設定
TITLE_FONT = 'Lato'
BODY_FONT = 'Lato'

# フォントサイズ定義（スライドからはみ出さないように調整）
TITLE_SIZE = Pt(42)
HEADING_SIZE = Pt(30)
SUBHEADING_SIZE = Pt(20)
BODY_SIZE = Pt(14)
CAPTION_SIZE = Pt(12)
TABLE_HEADER_SIZE = Pt(14)
TABLE_BODY_SIZE = Pt(12)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    total_slides = 10
    
    create_title_slide(prs, 1, total_slides)
    create_executive_summary(prs, 2, total_slides)
    create_current_analysis(prs, 3, total_slides)
    create_proposal(prs, 4, total_slides)
    create_schedule(prs, 5, total_slides)
    create_team_structure(prs, 6, total_slides)
    create_risk_management(prs, 7, total_slides)
    create_budget(prs, 8, total_slides)
    create_success_criteria(prs, 9, total_slides)
    create_conclusion(prs, 10, total_slides)
    
    prs.save('project_proposal.pptx')
    print("洗練されたプレゼンテーションが作成されました: project_proposal.pptx")

def apply_title_style(title_shape, text, font_size=TITLE_SIZE, color=ColorPalette.TEXT, align=PP_ALIGN.LEFT, bold=True):
    title_shape.text = text
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.alignment = align
    title_run = title_para.runs[0]
    title_run.font.name = TITLE_FONT
    title_run.font.size = font_size
    title_run.font.bold = bold
    title_run.font.color.rgb = color

def apply_body_style(body_shape, text_list, font_size=BODY_SIZE, color=ColorPalette.TEXT, para_spacing=Pt(8)):
    tf = body_shape.text_frame
    tf.clear()
    tf.word_wrap = True
    if not tf.paragraphs:
        tf.add_paragraph()
    p = tf.paragraphs[0]
    
    for i, original_text in enumerate(text_list):
        text_to_set = original_text
        if i > 0:
            p = tf.add_paragraph()
        if not original_text.strip():
            p.text = ""
            p.space_after = para_spacing
            continue
        level = 0
        space_before = Pt(0)
        is_heading = False
        if original_text.startswith('【') and original_text.endswith('】'):
            level = 0
            space_before = Pt(12) if i > 0 else Pt(0)
            is_heading = True
        elif original_text.startswith('• '):
            level = 1
            text_to_set = original_text[2:]
        else:
            level = 0
        
        p.text = text_to_set
        p.level = level
        p.space_after = para_spacing
        p.space_before = space_before
        
        if p.runs:
            run = p.runs[0]
            run.font.name = BODY_FONT
            run.font.size = font_size
            if is_heading:
                run.font.bold = True
                run.font.size = font_size + Pt(2)
                run.font.color.rgb = color
            else:
                run.font.bold = False
                run.font.color.rgb = color

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0.75), shadow=False, transparency=0, gradient_to=None, text=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        if gradient_to:
            shape.fill.gradient()
            shape.fill.gradient_stops[0].position = 0
            shape.fill.gradient_stops[0].color.rgb = fill_color
            shape.fill.gradient_stops[1].position = 1
            shape.fill.gradient_stops[1].color.rgb = gradient_to
            shape.fill.gradient_angle = 90
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            shape.fill.transparency = transparency
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    if shadow:
        shape.shadow.inherit = False
        shape.shadow.visible = True
        shape.shadow.blur_radius = Pt(3)
        shape.shadow.distance = Pt(2)
        shape.shadow.angle = 45
        try:
            shape.shadow.color.rgb = RGBColor(180, 180, 180)
            shape.shadow.transparency = 0.6
        except AttributeError:
            pass
    if text:
        shape.text = text
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = TITLE_FONT
        run.font.size = SUBHEADING_SIZE
        run.font.color.rgb = ColorPalette.HEADING_TEXT
    return shape

def add_background(slide, prs, type="solid", color=ColorPalette.BACKGROUND, gradient_to=None):
    background = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_color=color, line_color=None)
    return background

def add_footer(slide, prs, text="Your Company Name | Project Proposal", current_slide=1, total_slides=10):
    footer_shape = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), prs.slide_height - Inches(0.3), prs.slide_width, Inches(0.3), fill_color=ColorPalette.FOOTER_BG, line_color=None)
    footer_text = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.35), prs.slide_width - Inches(1.5), Inches(0.3))
    tf = footer_text.text_frame
    p = tf.paragraphs[0]
    p.text = f"{text} | {current_slide}/{total_slides}"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = BODY_FONT
    run.font.size = CAPTION_SIZE
    run.font.color.rgb = ColorPalette.FOOTER_TEXT
    return footer_shape

def create_table(slide, rows, cols, left, top, width, height):
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    return table

def set_table_cell_text(table, row, col, text, bold=False, alignment=PP_ALIGN.LEFT, font_size=None):
    cell = table.cell(row, col)
    para = cell.text_frame.paragraphs[0]
    para.text = text
    para.alignment = alignment
    if para.runs:
        run = para.runs[0]
        run.font.bold = bold
        if font_size:
            run.font.size = font_size

def create_title_slide(prs, current_slide, total_slides):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)
    header_line = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(6), Inches(0.05), fill_color=ColorPalette.ACCENT)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(10), Inches(2))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = "IT Development & System Implementation"
    title_p.alignment = PP_ALIGN.LEFT
    title_run = title_p.runs[0]
    title_run.font.name = TITLE_FONT
    title_run.font.size = TITLE_SIZE
    title_run.font.bold = True
    title_run.font.color.rgb = ColorPalette.TEXT
    subtitle_p = title_tf.add_paragraph()
    subtitle_p.text = "Project Proposal"
    subtitle_p.alignment = PP_ALIGN.LEFT
    subtitle_p.space_before = Pt(10)
    subtitle_run = subtitle_p.runs[0]
    subtitle_run.font.name = TITLE_FONT
    subtitle_run.font.size = SUBHEADING_SIZE
    subtitle_run.font.bold = False
    subtitle_run.font.color.rgb = ColorPalette.TEXT
    details_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(0.5))
    details_tf = details_box.text_frame
    details_p = details_tf.paragraphs[0]
    details_p.text = "March 30, 2025 | Your Company Name"
    details_p.alignment = PP_ALIGN.LEFT
    details_run = details_p.runs[0]
    details_run.font.name = BODY_FONT
    details_run.font.size = BODY_SIZE
    details_run.font.color.rgb = ColorPalette.FOOTER_TEXT
    add_footer(slide, prs, "Your Company Name | Project Proposal", current_slide, total_slides)

def create_executive_summary(prs, current_slide, total_slides):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)
    header = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1), fill_color=ColorPalette.HEADING_BG)
    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Executive Summary"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT
    left_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(5.5), Inches(4.5))
    left_content = [
        "【Project Objective】",
        "• Revamp the current business system to improve operational efficiency by 30%.",
        "• Build a foundation for digital transformation.",
        "",
        "【Key Proposal】",
        "• Implement a cloud-based integrated management system.",
        "• Utilize AI for business process automation and predictive analytics.",
        "",
        "【Expected Benefits】",
        "• Annual cost savings of ¥20 million.",
        "• Reduction in customer response time by 50%.",
        "• Enable data-driven decision-making.",
        "• Expand business opportunities through improved efficiency."
    ]
    apply_body_style(left_box, left_content, para_spacing=Pt(8))
    summary_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.5), Inches(5), Inches(1.2), fill_color=ColorPalette.LIGHT_ACCENT, line_color=ColorPalette.ACCENT, line_width=Pt(1))
    summary_text = slide.shapes.add_textbox(Inches(7.2), Inches(1.6), Inches(4.6), Inches(1))
    summary_tf = summary_text.text_frame
    summary_p = summary_tf.paragraphs[0]
    summary_p.text = "Key Project Information"
    summary_p.alignment = PP_ALIGN.CENTER
    summary_run = summary_p.runs[0]
    summary_run.font.name = TITLE_FONT
    summary_run.font.size = SUBHEADING_SIZE
    summary_run.font.bold = True
    summary_run.font.color.rgb = ColorPalette.TEXT
    info_box = slide.shapes.add_textbox(Inches(7), Inches(3), Inches(5), Inches(3))
    info_content = [
        "【Project Timeline】",
        "• Duration: 6 months (Apr 2025 - Sep 2025)",
        "",
        "【Budget Overview】",
        "• Initial investment: ¥35M",
        "• Annual operating cost: ¥8M",
        "",
        "【Return on Investment】",
        "• Payback period: 18 months",
        "• Efficiency gains: 30% in target processes"
    ]
    apply_body_style(info_box, info_content, para_spacing=Pt(8))
    add_footer(slide, prs, "Your Company Name | Project Proposal", current_slide, total_slides)

def create_current_analysis(prs, current_slide, total_slides):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)
    header = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1), fill_color=ColorPalette.HEADING_BG)
    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Current Situation & Challenges"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT
    table_width = Inches(11)
    table_height = Inches(5)
    table = create_table(slide, rows=6, cols=2, left=Inches(1.15), top=Inches(1.5), width=table_width, height=table_height)
    set_table_cell_text(table, 0, 0, "Current System Situation", bold=True, alignment=PP_ALIGN.CENTER)
    set_table_cell_text(table, 0, 1, "Key Challenges to Address", bold=True, alignment=PP_ALIGN.CENTER)
    current_situation = [
        "Core system operational for 8 years.",
        "Multiple systems lack integration, causing duplicate data entry.",
        "Increased maintenance costs due to legacy systems.",
        "Resource constraints in the on-premises environment.",
        "Lack of mobile support restricts remote work."
    ]
    challenges = [
        "Centralize data management and standardize business processes.",
        "Eliminate redundant work through automated system integration.",
        "Optimize costs by migrating to a cloud environment.",
        "Establish a remote work environment with mobile support.",
        "Enhance security and ensure compliance."
    ]
    for i, (current, challenge) in enumerate(zip(current_situation, challenges)):
        set_table_cell_text(table, i+1, 0, current)
        set_table_cell_text(table, i+1, 1, challenge)
    for col in table.columns:
        col.width = int(table_width / 2)
    add_footer(slide, prs, "Your Company Name | Project Proposal", current_slide, total_slides)

def create_proposal(prs, current_slide, total_slides):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)
    header = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1), fill_color=ColorPalette.HEADING_BG)
    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Proposal: Cloud Integrated Management System"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT
    features_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(11.3), Inches(2.5), fill_color=ColorPalette.LIGHT_ACCENT, line_color=ColorPalette.ACCENT, line_width=Pt(1))
    features_title = slide.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(10.9), Inches(0.5))
    features_tf = features_title.text_frame
    features_p = features_tf.paragraphs[0]
    features_p.text = "System Features"
    features_run = features_p.runs[0]
    features_run.font.name = TITLE_FONT
    features_run.font.size = SUBHEADING_SIZE
    features_run.font.bold = True
    features_run.font.color.rgb = ColorPalette.TEXT
    features_left = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(5.5), Inches(1.8))
    features_left_content = [
        "• Centralized management of all business data.",
        "• Cloud-based platform accessible from anywhere.",
        "• Intuitive user interface.",
        "• Real-time data synchronization and analysis."
    ]
    apply_body_style(features_left, features_left_content, para_spacing=Pt(6))
    features_right = slide.shapes.add_textbox(Inches(6.7), Inches(2.1), Inches(5.5), Inches(1.8))
    features_right_content = [
        "• Efficiency gains through business process automation.",
        "• AI-powered predictive analytics and decision support.",
        "• Flexible scalability and customization.",
        "• Enhanced security and compliance features."
    ]
    apply_body_style(features_right, features_right_content, para_spacing=Pt(6))
    functions_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.2), Inches(11.3), Inches(2.5), fill_color=ColorPalette.LIGHT_ACCENT, line_color=ColorPalette.ACCENT, line_width=Pt(1))
    functions_title = slide.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(10.9), Inches(0.5))
    functions_tf = functions_title.text_frame
    functions_p = functions_tf.paragraphs[0]
    functions_p.text = "Key Functions"
    functions_run = functions_p.runs[0]
    functions_run.font.name = TITLE_FONT
    functions_run.font.size = SUBHEADING_SIZE
    functions_run.font.bold = True
    functions_run.font.color.rgb = ColorPalette.TEXT
    table_width = Inches(10.9)
    table_height = Inches(1.6)
    functions_table = create_table(slide, rows=3, cols=2, left=Inches(1.2), top=Inches(4.8), width=table_width, height=table_height)
    key_functions = [
        ["Customer & Case Management", "Real-time Dashboards"],
        ["Workflow Automation", "Role-based Access Control"],
        ["Mobile Application Support", "API Integration Hub"]
    ]
    for row_idx, (left_func, right_func) in enumerate(key_functions):
        set_table_cell_text(functions_table, row_idx, 0, left_func)
        set_table_cell_text(functions_table, row_idx, 1, right_func)
    add_footer(slide, prs, "Your Company Name | Project Proposal", current_slide, total_slides)

def create_schedule(prs, current_slide, total_slides):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)
    header = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1), fill_color=ColorPalette.HEADING_BG)
    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Implementation Schedule (6-Month Plan)"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT
    table_width = Inches(11)
    table_height = Inches(5)
    schedule_table = create_table(slide, rows=5, cols=3, left=Inches(1.15), top=Inches(1.5), width=table_width, height=table_height)
    set_table_cell_text(schedule_table, 0, 0, "Phase", bold=True, alignment=PP_ALIGN.CENTER)
    set_table_cell_text(schedule_table, 0, 1, "Timeline", bold=True, alignment=PP_ALIGN.CENTER)
    set_table_cell_text(schedule_table, 0, 2, "Key Activities", bold=True, alignment=PP_ALIGN.CENTER)
    phases = [
        ["Phase 1:\nRequirements & Design", "Apr-May 2025", "• Detailed business requirement analysis.\n• System design and architecture finalization.\n• Data migration planning."],
        ["Phase 2:\nDevelopment & Build", "May-Jul 2025", "• Platform setup and core feature development.\n• Implementation of external system integrations.\n• User interface development."],
        ["Phase 3:\nTesting & Migration", "Jul-Aug 2025", "• Unit and integration testing.\n• User acceptance testing (UAT).\n• Data migration and system switchover preparation."],
        ["Phase 4:\nGo-Live & Stabilization", "Sep 2025", "• Phased production rollout.\n• User training sessions.\n• Establishment of operational support."]
    ]
    for row_idx, (phase, timeline, activities) in enumerate(phases):
        set_table_cell_text(schedule_table, row_idx + 1, 0, phase, bold=True)
        set_table_cell_text(schedule_table, row_idx + 1, 1, timeline, alignment=PP_ALIGN.CENTER)
        cell = schedule_table.cell(row_idx + 1, 2)
        tf_cell = cell.text_frame
        tf_cell.text = ""
        for i, line in enumerate(activities.split("\n")):
            p = tf_cell.add_paragraph() if i > 0 else tf_cell.paragraphs[0]
            p.text = line
            p.space_after = Pt(3)
            run = p.runs[0]
            run.font.name = BODY_FONT
            run.font.size = TABLE_BODY_SIZE
            run.font.color.rgb = ColorPalette.TEXT
    schedule_table.columns[0].width = Inches(2.5)
    schedule_table.columns[1].width = Inches(2)
    schedule_table.columns[2].width = Inches(6.5)
    add_footer(slide, prs, "Your Company Name | Project Proposal", current_slide, total_slides)

def create_team_structure(prs, current_slide, total_slides):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)
    header = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1), fill_color=ColorPalette.HEADING_BG)
    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Project Team Structure"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT
    sponsor_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.15), Inches(1.7), Inches(3), Inches(1), fill_color=ColorPalette.HEADING_BG, line_color=ColorPalette.ACCENT, line_width=Pt(2))
    sponsor_text = slide.shapes.add_textbox(Inches(5.25), Inches(1.9), Inches(2.8), Inches(0.7))
    sponsor_tf = sponsor_text.text_frame
    sponsor_p = sponsor_tf.paragraphs[0]
    sponsor_p.text = "Project Sponsor:\nHead of Corporate Planning"
    sponsor_p.alignment = PP_ALIGN.CENTER
    sponsor_run = sponsor_p.runs[0]
    sponsor_run.font.name = BODY_FONT
    sponsor_run.font.size = BODY_SIZE
    sponsor_run.font.bold = True
    sponsor_run.font.color.rgb = ColorPalette.HEADING_TEXT
    pm_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.15), Inches(3.1), Inches(3), Inches(1), fill_color=ColorPalette.ACCENT, line_color=None)
    pm_text = slide.shapes.add_textbox(Inches(5.25), Inches(3.3), Inches(2.8), Inches(0.7))
    pm_tf = pm_text.text_frame
    pm_p = pm_tf.paragraphs[0]
    pm_p.text = "Project Manager:\nIT Department Manager"
    pm_p.alignment = PP_ALIGN.CENTER
    pm_run = pm_p.runs[0]
    pm_run.font.name = BODY_FONT
    pm_run.font.size = BODY_SIZE
    pm_run.font.bold = True
    pm_run.font.color.rgb = ColorPalette.HEADING_TEXT
    arrow1 = add_shape(slide, MSO_SHAPE.DOWN_ARROW, Inches(6.4), Inches(2.72), Inches(0.5), Inches(0.35), fill_color=ColorPalette.ACCENT)
    tech_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(4.5), Inches(3), Inches(1), fill_color=ColorPalette.LIGHT_ACCENT, line_color=ColorPalette.ACCENT, line_width=Pt(1))
    tech_text = slide.shapes.add_textbox(Inches(1.6), Inches(4.7), Inches(2.8), Inches(0.7))
    tech_tf = tech_text.text_frame
    tech_p = tech_tf.paragraphs[0]
    tech_p.text = "Technical Lead:\nLead Systems Developer"
    tech_p.alignment = PP_ALIGN.CENTER
    tech_run = tech_p.runs[0]
    tech_run.font.name = BODY_FONT
    tech_run.font.size = BODY_SIZE
    tech_run.font.bold = True
    tech_run.font.color.rgb = ColorPalette.TEXT
    biz_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(4.5), Inches(3), Inches(1), fill_color=ColorPalette.LIGHT_ACCENT, line_color=ColorPalette.ACCENT, line_width=Pt(1))
    biz_text = slide.shapes.add_textbox(Inches(8.9), Inches(4.7), Inches(2.8), Inches(0.7))
    biz_tf = biz_text.text_frame
    biz_p = biz_tf.paragraphs[0]
    biz_p.text = "Business Process Owners:\nRepresentatives from each dept."
    biz_p.alignment = PP_ALIGN.CENTER
    biz_run = biz_p.runs[0]
    biz_run.font.name = BODY_FONT
    biz_run.font.size = BODY_SIZE
    biz_run.font.bold = True
    biz_run.font.color.rgb = ColorPalette.TEXT
    arrow2 = add_shape(slide, MSO_SHAPE.LEFT_UP_ARROW, Inches(5.15), Inches(4), Inches(1), Inches(0.7), fill_color=ColorPalette.ACCENT)
    arrow3 = add_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(7.15), Inches(4), Inches(1), Inches(0.7), fill_color=ColorPalette.ACCENT)
    comm_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.65), Inches(5.8), Inches(6), Inches(1), fill_color=ColorPalette.LIGHT_ACCENT, line_color=ColorPalette.ACCENT, line_width=Pt(1))
    comm_text = slide.shapes.add_textbox(Inches(3.75), Inches(5.9), Inches(5.8), Inches(0.9))
    comm_tf = comm_text.text_frame
    comm_p = comm_tf.paragraphs[0]
    comm_p.text = "Communication Plan: Weekly meetings (online), Monthly steering committee (in-person), Daily stand-ups for development team"
    comm_p.alignment = PP_ALIGN.CENTER
    comm_run = comm_p.runs[0]
    comm_run.font.name = BODY_FONT
    comm_run.font.size = BODY_SIZE
    comm_run.font.color.rgb = ColorPalette.TEXT
    add_footer(slide, prs, "Your Company Name | Project Proposal", current_slide, total_slides)

def create_risk_management(prs, current_slide, total_slides):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)
    header = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1), fill_color=ColorPalette.HEADING_BG)
    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Risk Management Plan"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT
    subtitle_box = slide.shapes.add_textbox(Inches(1.15), Inches(1.3), Inches(11), Inches(0.5))
    subtitle_tf = subtitle_box.text_frame
    subtitle_p = subtitle_tf.paragraphs[0]
    subtitle_p.text = "Key Risks & Mitigation Strategies"
    subtitle_run = subtitle_p.runs[0]
    subtitle_run.font.name = TITLE_FONT
    subtitle_run.font.size = SUBHEADING_SIZE
    subtitle_run.font.bold = True
    subtitle_run.font.color.rgb = ColorPalette.TEXT
    table_width = Inches(11)
    table_height = Inches(4.5)
    risk_table = create_table(slide, rows=6, cols=2, left=Inches(1.15), top=Inches(1.9), width=table_width, height=table_height)
    set_table_cell_text(risk_table, 0, 0, "Risk", bold=True, alignment=PP_ALIGN.CENTER)
    set_table_cell_text(risk_table, 0, 1, "Mitigation Strategy", bold=True, alignment=PP_ALIGN.CENTER)
    risks = [
        ["Scope Creep / Changes Leading to Delays", "Agile methodology, regular requirement reviews, strict change control."],
        ["Data Loss / Inconsistency During Migration", "Pre-migration data cleansing, phased approach, dual validation."],
        ["Low User Adoption", "Early user involvement, comprehensive training, continuous feedback loop."],
        ["Integration Issues with Existing Systems", "Detailed interface design, phased integration testing, fallback mechanisms."],
        ["Security Incidents", "Security design reviews, vulnerability assessments, incident response plan."]
    ]
    for row_idx, (risk, mitigation) in enumerate(risks):
        set_table_cell_text(risk_table, row_idx + 1, 0, risk, bold=True)
        set_table_cell_text(risk_table, row_idx + 1, 1, mitigation)
    risk_table.columns[0].width = Inches(4)
    risk_table.columns[1].width = Inches(7)
    add_footer(slide, prs, "Your Company Name | Project Proposal", current_slide, total_slides)

def create_budget(prs, current_slide, total_slides):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)
    header = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1), fill_color=ColorPalette.HEADING_BG)
    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Budget Plan & ROI"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT
    subtitle1 = slide.shapes.add_textbox(Inches(1.15), Inches(1.3), Inches(5.3), Inches(0.4))
    subtitle1_tf = subtitle1.text_frame
    subtitle1_p = subtitle1_tf.paragraphs[0]
    subtitle1_p.text = "Initial Investment"
    subtitle1_run = subtitle1_p.runs[0]
    subtitle1_run.font.name = TITLE_FONT
    subtitle1_run.font.size = SUBHEADING_SIZE
    subtitle1_run.font.bold = True
    subtitle1_run.font.color.rgb = ColorPalette.TEXT
    initial_table = create_table(slide, rows=6, cols=2, left=Inches(1.15), top=Inches(1.8), width=Inches(5.3), height=Inches(2.5))
    initial_items = [
        ["Item", "Cost"],
        ["Design & Development", "¥20M"],
        ["Hardware & Cloud Setup", "¥5M"],
        ["Data Migration & Testing", "¥6M"],
        ["Training & Support", "¥4M"],
        ["Total Initial Cost", "¥35M"]
    ]
    for row_idx, (item, cost) in enumerate(initial_items):
        set_table_cell_text(initial_table, row_idx, 0, item, bold=(row_idx==0 or row_idx==5))
        set_table_cell_text(initial_table, row_idx, 1, cost, bold=(row_idx==0 or row_idx==5), alignment=PP_ALIGN.RIGHT)
    subtitle2 = slide.shapes.add_textbox(Inches(1.15), Inches(4.4), Inches(5.3), Inches(0.4))
    subtitle2_tf = subtitle2.text_frame
    subtitle2_p = subtitle2_tf.paragraphs[0]
    subtitle2_p.text = "Annual Running Costs"
    subtitle2_run = subtitle2_p.runs[0]
    subtitle2_run.font.name = TITLE_FONT
    subtitle2_run.font.size = SUBHEADING_SIZE
    subtitle2_run.font.bold = True
    subtitle2_run.font.color.rgb = ColorPalette.TEXT
    running_table = create_table(slide, rows=5, cols=2, left=Inches(1.15), top=Inches(4.9), width=Inches(5.3), height=Inches(1.9))
    running_items = [
        ["Item", "Cost"],
        ["Cloud Infrastructure", "¥3M"],
        ["Licensing Fees", "¥2M"],
        ["Maintenance & Support", "¥3M"],
        ["Total Annual Cost", "¥8M"]
    ]
    for row_idx, (item, cost) in enumerate(running_items):
        set_table_cell_text(running_table, row_idx, 0, item, bold=(row_idx==0 or row_idx==4))
        set_table_cell_text(running_table, row_idx, 1, cost, bold=(row_idx==0 or row_idx==4), alignment=PP_ALIGN.RIGHT)
    roi_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.95), Inches(1.8), Inches(5.3), Inches(5), fill_color=ColorPalette.LIGHT_ACCENT, line_color=ColorPalette.ACCENT, line_width=Pt(1))
    roi_title = slide.shapes.add_textbox(Inches(7.15), Inches(1.9), Inches(5), Inches(0.4))
    roi_tf = roi_title.text_frame
    roi_p = roi_tf.paragraphs[0]
    roi_p.text = "Return on Investment (ROI)"
    roi_p.alignment = PP_ALIGN.CENTER
    roi_run = roi_p.runs[0]
    roi_run.font.name = TITLE_FONT
    roi_run.font.size = SUBHEADING_SIZE
    roi_run.font.bold = True
    roi_run.font.color.rgb = ColorPalette.TEXT
    roi_content = slide.shapes.add_textbox(Inches(7.15), Inches(2.5), Inches(5), Inches(4))
    roi_text = [
        "【Cost Savings】",
        "• Labor cost reduction (efficiency): ¥12M/year",
        "• System consolidation savings: ¥8M/year",
        "• Total annual savings: ¥20M/year",
        "",
        "【Qualitative Benefits】",
        "• Faster decision-making",
        "• Improved customer satisfaction",
        "• Strategic advantage through data utilization",
        "",
        "【Payback Period】",
        "• Initial investment: ¥35M",
        "• Annual savings: ¥20M",
        "• Payback period: ~18 months"
    ]
    apply_body_style(roi_content, roi_text, para_spacing=Pt(8))
    add_footer(slide, prs, "Your Company Name | Project Proposal", current_slide, total_slides)

def create_success_criteria(prs, current_slide, total_slides):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)
    header = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1), fill_color=ColorPalette.HEADING_BG)
    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Success Criteria & Evaluation"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT
    subtitle1 = slide.shapes.add_textbox(Inches(1.15), Inches(1.3), Inches(11), Inches(0.4))
    subtitle1_tf = subtitle1.text_frame
    subtitle1_p = subtitle1_tf.paragraphs[0]
    subtitle1_p.text = "Key Performance Indicators (KPIs)"
    subtitle1_run = subtitle1_p.runs[0]
    subtitle1_run.font.name = TITLE_FONT
    subtitle1_run.font.size = SUBHEADING_SIZE
    subtitle1_run.font.bold = True
    subtitle1_run.font.color.rgb = ColorPalette.TEXT
    kpi_table = create_table(slide, rows=5, cols=4, left=Inches(1.15), top=Inches(1.8), width=Inches(11), height=Inches(2.5))
    set_table_cell_text(kpi_table, 0, 0, "System Performance Metric", bold=True, alignment=PP_ALIGN.CENTER)
    set_table_cell_text(kpi_table, 0, 1, "Target", bold=True, alignment=PP_ALIGN.CENTER)
    set_table_cell_text(kpi_table, 0, 2, "Business Impact Metric", bold=True, alignment=PP_ALIGN.CENTER)
    set_table_cell_text(kpi_table, 0, 3, "Target", bold=True, alignment=PP_ALIGN.CENTER)
    system_metrics = [
        ["Response Time", "< 2 seconds (peak)"],
        ["Availability", "> 99.9%"],
        ["Concurrent Users", "Up to 300"],
        ["Backup Recovery Time", "< 4 hours"]
    ]
    business_metrics = [
        ["Process Time Reduction", "30%"],
        ["Customer Response Time", "50% improvement"],
        ["Data Entry Error Reduction", "90%"],
        ["User Satisfaction", "> 80%"]
    ]
    for row_idx in range(4):
        set_table_cell_text(kpi_table, row_idx + 1, 0, system_metrics[row_idx][0])
        set_table_cell_text(kpi_table, row_idx + 1, 1, system_metrics[row_idx][1], alignment=PP_ALIGN.CENTER)
        set_table_cell_text(kpi_table, row_idx + 1, 2, business_metrics[row_idx][0])
        set_table_cell_text(kpi_table, row_idx + 1, 3, business_metrics[row_idx][1], alignment=PP_ALIGN.CENTER)
    kpi_table.columns[0].width = Inches(3.5)
    kpi_table.columns[1].width = Inches(2)
    kpi_table.columns[2].width = Inches(3.5)
    kpi_table.columns[3].width = Inches(2)
    subtitle2 = slide.shapes.add_textbox(Inches(1.15), Inches(4.5), Inches(11), Inches(0.4))
    subtitle2_tf = subtitle2.text_frame
    subtitle2_p = subtitle2_tf.paragraphs[0]
    subtitle2_p.text = "Evaluation Method"
    subtitle2_run = subtitle2_p.runs[0]
    subtitle2_run.font.name = TITLE_FONT
    subtitle2_run.font.size = SUBHEADING_SIZE
    subtitle2_run.font.bold = True
    subtitle2_run.font.color.rgb = ColorPalette.TEXT
    criteria_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(2))
    criteria_points = [
        "• Quarterly performance measurement reports.",
        "• Monthly user satisfaction surveys.",
        "• Regular tracking of business efficiency metrics.",
        "• Continuous monitoring via real-time dashboards."
    ]
    apply_body_style(criteria_box, criteria_points, para_spacing=Pt(8))
    add_footer(slide, prs, "Your Company Name | Project Proposal", current_slide, total_slides)

def create_conclusion(prs, current_slide, total_slides):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs, color=ColorPalette.HEADING_BG)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "Conclusion & Next Steps"
    title_p.alignment = PP_ALIGN.LEFT
    title_run = title_p.runs[0]
    title_run.font.name = TITLE_FONT
    title_run.font.size = HEADING_SIZE
    title_run.font.bold = True
    title_run.font.color.rgb = ColorPalette.HEADING_TEXT
    summary_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(1.5))
    summary_text = [
        "• Implement cloud-based system for 30% efficiency gain.",
        "• Phased 6-month rollout minimizes business disruption.",
        "• Investment: ¥35M initial, ¥8M annual. ROI within 18 months."
    ]
    apply_body_style(summary_box, summary_text, color=ColorPalette.HEADING_TEXT, para_spacing=Pt(8))
    next_steps_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1.5))
    next_text = [
        "• Final review and approval of proposal (within 1 week).",
        "• Project kick-off meeting (within 2 weeks of approval).",
        "• Commence detailed requirements definition (First week of April)."
    ]
    apply_body_style(next_steps_box, next_text, color=ColorPalette.HEADING_TEXT, para_spacing=Pt(8))
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
    contact_tf = contact_box.text_frame
    contact_p = contact_tf.paragraphs[0]
    contact_p.text = "Contact: Taro Yamada | yamada.taro@example.com | 03-1234-5678"
    contact_p.alignment = PP_ALIGN.LEFT
    contact_run = contact_p.runs[0]
    contact_run.font.name = BODY_FONT
    contact_run.font.size = BODY_SIZE
    contact_run.font.color.rgb = ColorPalette.FOOTER_BG
    add_footer(slide, prs, "Your Company Name | Project Proposal", current_slide, total_slides)

if __name__ == "__main__":
    create_presentation()