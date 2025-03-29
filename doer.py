import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL

# 白と黒を基調としたシンプルなカラーパレット
class ColorPalette:
    BACKGROUND = RGBColor(255, 255, 255)  # ホワイト
    TEXT = RGBColor(0, 0, 0)             # ブラック
    ACCENT = RGBColor(50, 50, 50)         # ダークグレー
    FOOTER_BG = RGBColor(240, 240, 240)   # ライトグレー
    FOOTER_TEXT = RGBColor(100, 100, 100) # ミディアムグレー
    HEADING_BG = RGBColor(0, 0, 0)        # ブラック
    HEADING_TEXT = RGBColor(255, 255, 255) # ホワイト
    PANEL_BG = RGBColor(255, 255, 255)   # ホワイト
    PANEL_BORDER = RGBColor(220, 220, 220) # ライトグレーボーダー

# モダンで洗練されたフォント設定
TITLE_FONT = 'Lato'  # Google Fonts でも利用可能なモダンフォント
BODY_FONT = 'Lato'
# ALT_FONT = 'Helvetica Neue' # 代替はコメントアウト

# フォントサイズ定義 (余白を意識して調整)
TITLE_SIZE = Pt(48)
HEADING_SIZE = Pt(32)
SUBHEADING_SIZE = Pt(20)
BODY_SIZE = Pt(16)
CAPTION_SIZE = Pt(12)

def create_presentation():
    prs = Presentation()
    
    # スライドサイズを16:9に設定
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # スライドを作成
    create_title_slide(prs)
    create_executive_summary(prs)
    create_current_analysis(prs)
    create_proposal(prs)
    create_schedule(prs)
    create_team_structure(prs)
    create_risk_management(prs)
    create_budget(prs)
    create_success_criteria(prs)
    create_conclusion(prs)
    
    # プレゼンテーションを保存
    prs.save('great1.pptx')
    print("洗練されたプレゼンテーションが作成されました: great1.pptx")

def apply_title_style(title_shape, text, font_size=TITLE_SIZE, color=ColorPalette.TEXT, align=PP_ALIGN.LEFT, bold=True):
    """タイトルのスタイルを適用する"""
    title_shape.text = text
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.alignment = align
    title_run = title_para.runs[0]
    title_run.font.name = TITLE_FONT
    title_run.font.size = font_size
    title_run.font.bold = bold
    title_run.font.color.rgb = color

def apply_body_style(body_shape, text_list, font_size=BODY_SIZE, color=ColorPalette.TEXT, para_spacing=Pt(12)):
    """本文のスタイルを適用する"""
    tf = body_shape.text_frame
    tf.clear()

    # 最初の段落を確保
    if not tf.paragraphs:
        tf.add_paragraph()
    p = tf.paragraphs[0]

    for i, original_text in enumerate(text_list):
        text_to_set = original_text  # スタイル適用用のテキスト

        # 2番目以降の要素のために新しい段落を追加
        if i > 0:
            p = tf.add_paragraph()

        # 空文字列または空白のみの文字列を処理
        if not original_text.strip():
            p.text = ""  # 空のテキストを設定
            p.space_after = para_spacing
            continue  # 次の要素へ

        # レベル決定とテキスト変更（割り当て前）
        level = 0
        space_before = Pt(0)
        is_heading = False
        is_bullet1 = False
        is_bullet2 = False

        if original_text.startswith('【') and original_text.endswith('】'):
            level = 0
            space_before = Pt(15) if i > 0 else Pt(0)
            is_heading = True
        elif original_text.startswith('• '):
            level = 1
            text_to_set = original_text[2:]
            is_bullet1 = True
        elif original_text.startswith('  - '):
            level = 2
            text_to_set = original_text[4:]
            is_bullet2 = True
        else:
            level = 0

        # テキスト、レベル、スペースを設定
        p.text = text_to_set
        p.level = level
        p.space_after = para_spacing
        p.space_before = space_before

        # runが存在するか確認してスタイルを適用
        if p.runs:
            run = p.runs[0]
            run.font.name = BODY_FONT
            run.font.size = font_size

            # 強調表示などの特定のスタイル
            if is_heading:
                run.font.bold = True
                run.font.size = font_size + Pt(4) # 見出しを少し大きく
                run.font.color.rgb = color # 同じテキストカラーを使用
            else:
                run.font.bold = False
                run.font.color.rgb = color
        # else:
            # テキストが空でないのにrunが作成されない場合（通常は発生しない）
            # print(f"警告: テキスト '{text_to_set}' のrunが作成されませんでした")
            # pass

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, 
              line_width=Pt(0.75), shadow=False, transparency=0, gradient_to=None, text=None):
    """洗練された図形を追加する (影はデフォルトOFF)"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    
    # 塗りつぶし設定
    if fill_color:
        if gradient_to:
            # グラデーション
            shape.fill.gradient()
            shape.fill.gradient_stops[0].position = 0
            shape.fill.gradient_stops[0].color.rgb = fill_color
            shape.fill.gradient_stops[1].position = 1
            shape.fill.gradient_stops[1].color.rgb = gradient_to
            shape.fill.gradient_angle = 90  # 左から右へのグラデーション
        else:
            # 単色
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            shape.fill.transparency = transparency
    else:
        # 塗りつぶしなし
        shape.fill.background()
    
    # 線の設定
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()  # 線なし
    
    # 影の設定
    if shadow:
        shape.shadow.inherit = False
        shape.shadow.visible = True
        shape.shadow.blur_radius = Pt(5)
        shape.shadow.distance = Pt(3)
        shape.shadow.angle = 45
        try:
            shape.shadow.color.rgb = RGBColor(180, 180, 180) # 薄いグレーの影
            shape.shadow.transparency = 0.5
        except AttributeError:
            pass
    
    # テキストがある場合は追加
    if text:
        shape.text = text
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = TITLE_FONT # タイトルと同じフォント
        run.font.size = SUBHEADING_SIZE
        run.font.color.rgb = ColorPalette.HEADING_TEXT # 白抜き文字
    
    return shape

def add_background(slide, prs, type="solid", color=ColorPalette.BACKGROUND, gradient_to=None):
    """スライドの背景を設定する (シンプルに単色のみ)"""
    background = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height,
        fill_color=color,
        line_color=None # 背景に線は不要
    )
    return background

def add_footer(slide, prs, text="Your Company Name | Project Proposal | 2025", include_page_number=True, page_num=None):
    """フッターを追加する (シンプルデザイン)"""
    footer_shape = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), prs.slide_height - Inches(0.3),
        prs.slide_width, Inches(0.3),
        fill_color=ColorPalette.FOOTER_BG,
        line_color=None # フッター背景に線は不要
    )
    
    # フッターテキスト
    footer_text = slide.shapes.add_textbox(
        Inches(0.5), prs.slide_height - Inches(0.35),
        prs.slide_width - Inches(1.5), Inches(0.3)
    )
    tf = footer_text.text_frame
    p = tf.paragraphs[0]
    
    if include_page_number and page_num:
        p.text = f"{text} | {page_num}"
    else:
        p.text = text
    
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = BODY_FONT
    run.font.size = CAPTION_SIZE
    run.font.color.rgb = ColorPalette.FOOTER_TEXT
    
    return footer_shape

def create_title_slide(prs):
    """洗練された表紙スライドの作成 (ミニマルデザイン)"""
    slide_layout = prs.slide_layouts[5] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色
    add_background(slide, prs)

    # タイトル
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(11), Inches(2)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = "IT Development & System Implementation"
    title_p.alignment = PP_ALIGN.LEFT
    title_run = title_p.runs[0]
    title_run.font.name = TITLE_FONT
    title_run.font.size = TITLE_SIZE
    title_run.font.bold = True
    title_run.font.color.rgb = ColorPalette.TEXT

    # サブタイトル
    subtitle_p = title_tf.add_paragraph()
    subtitle_p.text = "Project Proposal"
    subtitle_p.alignment = PP_ALIGN.LEFT
    subtitle_run = subtitle_p.runs[0]
    subtitle_run.font.name = TITLE_FONT
    subtitle_run.font.size = SUBHEADING_SIZE
    subtitle_run.font.bold = False
    subtitle_run.font.color.rgb = ColorPalette.TEXT

    # 日付と会社名 (フッターに含めるため、ここでは削除またはシンプルに)
    details_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.5),
        Inches(11), Inches(0.5)
    )
    details_tf = details_box.text_frame
    details_p = details_tf.paragraphs[0]
    details_p.text = "March 30, 2025 | Your Company Name"
    details_p.alignment = PP_ALIGN.LEFT
    details_run = details_p.runs[0]
    details_run.font.name = BODY_FONT
    details_run.font.size = BODY_SIZE
    details_run.font.color.rgb = ColorPalette.FOOTER_TEXT # 少し薄い色

    # フッター
    add_footer(slide, prs, text="Your Company Name", page_num="1/10") # ページ番号のみ表示するなど調整

def create_executive_summary(prs):
    """洗練されたエグゼクティブサマリーのスライド"""
    slide_layout = prs.slide_layouts[5] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs)
    
    # ヘッダーバー (シンプルに)
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1),
        fill_color=ColorPalette.HEADING_BG
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(12), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Executive Summary"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT

    # コンテンツエリア (余白を多く取る)
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5),
        Inches(11), Inches(5)
    )
    
    summary_points = [
        "【Project Objective】",
        "• Revamp the current business system to improve operational efficiency by 30%.",
        "• Build a foundation for digital transformation.",
        "",
        "【Key Proposal】",
        "• Implement a cloud-based integrated management system.",
        "• Utilize AI for business process automation and predictive analytics.",
        "",
        "【Expected Benefits】",
        "• Achieve annual cost savings of ¥20 million and reduce customer response time by 50%.",
        "• Enable data-driven decision-making and expand business opportunities."
    ]
    apply_body_style(content_box, summary_points, para_spacing=Pt(14))

    # 予算・期間の情報 (シンプルにテキストで)
    info_box = slide.shapes.add_textbox(
        Inches(8), Inches(5.5),
        Inches(4.5), Inches(1.2)
    )
    info_points = [
        "• Duration: 6 months (Apr 2025 - Sep 2025)",
        "• Budget: ¥35M initial, ¥8M annual running cost",
        "• ROI: Estimated payback in 18 months"
    ]
    # Apply style with smaller font and lighter color
    tf = info_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    for i, text in enumerate(info_points):
        if i > 0:
            p = tf.add_paragraph()
        p.text = text
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.name = BODY_FONT
        run.font.size = Pt(14)
        run.font.color.rgb = ColorPalette.FOOTER_TEXT

    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="2/10")

def create_current_analysis(prs):
    """洗練された現状分析と課題のスライド"""
    slide_layout = prs.slide_layouts[5] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1),
        fill_color=ColorPalette.HEADING_BG
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(12), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Current Situation & Challenges"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT

    # 左側: 現状 (シンプルにテキストボックス)
    current_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5),
        Inches(5.5), Inches(5)
    )
    current_tf = current_box.text_frame
    current_p = current_tf.paragraphs[0]
    current_p.text = "Current System Situation"
    current_run = current_p.runs[0]
    current_run.font.name = TITLE_FONT
    current_run.font.size = SUBHEADING_SIZE
    current_run.font.bold = True
    current_run.font.color.rgb = ColorPalette.TEXT
    current_p.space_after = Pt(12)

    current_state = [
        "• Core system operational for 8 years.",
        "• Multiple systems lack integration, causing duplicate data entry.",
        "• Increased maintenance costs due to legacy systems.",
        "• Resource constraints in the on-premises environment.",
        "• Lack of mobile support restricts remote work."
    ]
    apply_body_style(current_box, current_state, font_size=BODY_SIZE, para_spacing=Pt(12))
    # Remove the first paragraph which was the subheading
    current_tf.paragraphs[0].clear()

    # 右側: 課題 (シンプルにテキストボックス)
    challenge_box = slide.shapes.add_textbox(
        Inches(7), Inches(1.5),
        Inches(5.5), Inches(5)
    )
    challenge_tf = challenge_box.text_frame
    challenge_p = challenge_tf.paragraphs[0]
    challenge_p.text = "Key Challenges to Address"
    challenge_run = challenge_p.runs[0]
    challenge_run.font.name = TITLE_FONT
    challenge_run.font.size = SUBHEADING_SIZE
    challenge_run.font.bold = True
    challenge_run.font.color.rgb = ColorPalette.TEXT
    challenge_p.space_after = Pt(12)

    challenges = [
        "• Centralize data management and standardize business processes.",
        "• Eliminate redundant work through automated system integration.",
        "• Optimize costs by migrating to a cloud environment.",
        "• Establish a remote work environment with mobile support.",
        "• Enhance security and ensure compliance."
    ]
    apply_body_style(challenge_box, challenges, font_size=BODY_SIZE, para_spacing=Pt(12))
    challenge_tf.paragraphs[0].clear()

    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="3/10")

def create_proposal(prs):
    """提案内容のスライド (ミニマル)"""
    slide_layout = prs.slide_layouts[5] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1),
        fill_color=ColorPalette.HEADING_BG
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(12), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Proposal: Cloud Integrated Management System"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT

    # メインコンテンツエリア
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5),
        Inches(11), Inches(5.5)
    )
    
    main_content = [
        "【System Features】",
        "• Centralized management of all business data.",
        "• Cloud-based platform accessible from anywhere.",
        "• Intuitive user interface.",
        "• Real-time data synchronization and analysis.",
        "• Efficiency gains through business process automation.",
        "• AI-powered predictive analytics and decision support.",
        "• Flexible scalability and customization.",
        "",
        "【Key Functions】",
        "• Integrated customer information and case management.",
        "• Real-time dashboards.",
        "• Workflow automation.",
        "• Role-based access control and security.",
        "• Mobile application support."
    ]
    apply_body_style(content_box, main_content, para_spacing=Pt(14))

    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="4/10")

def create_schedule(prs):
    """導入スケジュールのスライド (テキストベース)"""
    slide_layout = prs.slide_layouts[5] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1),
        fill_color=ColorPalette.HEADING_BG
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(12), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Implementation Schedule (6-Month Plan)"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT

    # スケジュールコンテンツ
    schedule_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5),
        Inches(11), Inches(5.5)
    )
    
    schedule_text = [
        "【Phase 1: Requirements & Design (Apr-May)】",
        "• Detailed business requirement analysis.",
        "• System design and architecture finalization.",
        "• Data migration planning.",
        "",
        "【Phase 2: Development & Build (May-Jul)】",
        "• Platform setup and core feature development.",
        "• Implementation of external system integrations.",
        "• User interface development.",
        "",
        "【Phase 3: Testing & Migration (Jul-Aug)】",
        "• Unit and integration testing.",
        "• User acceptance testing (UAT).",
        "• Data migration and system switchover preparation.",
        "",
        "【Phase 4: Go-Live & Stabilization (Sep)】",
        "• Phased production rollout.",
        "• User training sessions.",
        "• Establishment of operational support."
    ]
    apply_body_style(schedule_box, schedule_text, font_size=BODY_SIZE, para_spacing=Pt(12))

    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="5/10")

def create_team_structure(prs):
    """実施体制のスライド (シンプル)"""
    slide_layout = prs.slide_layouts[5] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1),
        fill_color=ColorPalette.HEADING_BG
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(12), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Project Team Structure"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT

    # 体制コンテンツ (2カラム風)
    # 左カラム: 推進体制
    team_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5),
        Inches(5.5), Inches(5)
    )
    team_tf = team_box.text_frame
    team_p = team_tf.paragraphs[0]
    team_p.text = "Core Project Team"
    team_run = team_p.runs[0]
    team_run.font.name = TITLE_FONT
    team_run.font.size = SUBHEADING_SIZE
    team_run.font.bold = True
    team_p.space_after = Pt(12)

    team_text = [
        "• Project Sponsor: Head of Corporate Planning",
        "• Project Manager: IT Department Manager",
        "• Technical Lead: Lead Systems Developer",
        "• Business Process Owners: Representatives from each dept."
    ]
    apply_body_style(team_box, team_text, para_spacing=Pt(12))
    team_tf.paragraphs[0].clear()

    # 右カラム: 役割とコミュニケーション
    role_comm_box = slide.shapes.add_textbox(
        Inches(7), Inches(1.5),
        Inches(5.5), Inches(5)
    )
    role_comm_tf = role_comm_box.text_frame
    role_comm_tf.clear() # Ensure the text frame is empty before adding

    # 役割 ヘッダー追加
    role_p = role_comm_tf.paragraphs[0]
    role_p.text = "Roles & Responsibilities"
    role_run = role_p.runs[0]
    role_run.font.name = TITLE_FONT
    role_run.font.size = SUBHEADING_SIZE
    role_run.font.bold = True
    role_p.space_after = Pt(12)

    # 役割 テキスト追加 (apply_body_style を使用しない直接追加に変更)
    role_text = [
        "• Req. & Design: Our Consultants + Client Business Experts",
        "• Development: Our Engineering Team (5 members)",
        "• Testing & QA: Our QA Team + Client Testers",
        "• Implementation & Training: Our Support Team"
    ]
    for text in role_text:
        p = role_comm_tf.add_paragraph()
        if text.startswith("• "):
            p.text = text[2:]
            p.level = 1
        else:
            p.text = text
            p.level = 0
        p.space_after = Pt(12)
        run = p.runs[0]
        run.font.name = BODY_FONT
        run.font.size = BODY_SIZE
        run.font.color.rgb = ColorPalette.TEXT

    # コミュニケーション ヘッダー追加
    comm_p = role_comm_tf.add_paragraph()
    comm_p.text = "Communication Plan"
    comm_run = comm_p.runs[0]
    comm_run.font.name = TITLE_FONT
    comm_run.font.size = SUBHEADING_SIZE
    comm_run.font.bold = True
    comm_p.space_before = Pt(18) # Add space before the heading
    comm_p.space_after = Pt(12)

    # コミュニケーション テキスト追加
    comm_text = [
        "• Weekly Progress Meetings (Online)",
        "• Monthly Steering Committee (In-person)",
        "• Daily Stand-ups (Development Team)",
        "• Shared Issue Tracking System"
    ]
    for text in comm_text:
        p = role_comm_tf.add_paragraph()
        if text.startswith("• "):
            p.text = text[2:]
            p.level = 1
        else:
            p.text = text
            p.level = 0
        p.level = 1 # Indent communication points
        p.space_after = Pt(12)
        run = p.runs[0]
        run.font.name = BODY_FONT
        run.font.size = BODY_SIZE
        run.font.color.rgb = ColorPalette.TEXT

    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="6/10")

def create_risk_management(prs):
    """リスク管理計画のスライド (シンプル)"""
    slide_layout = prs.slide_layouts[5] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1),
        fill_color=ColorPalette.HEADING_BG
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(12), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Risk Management Plan"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT

    # リスクコンテンツ
    risk_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5),
        Inches(11), Inches(5.5)
    )
    
    risk_points = [
        "【Key Risks & Mitigation Strategies】",
        "",
        "1. Risk: Scope Creep / Changes Leading to Delays",
           "   Mitigation: Agile methodology, regular requirement reviews, strict change control.",
        "",
        "2. Risk: Data Loss / Inconsistency During Migration",
           "   Mitigation: Pre-migration data cleansing, phased approach, dual validation.",
        "",
        "3. Risk: Low User Adoption",
           "   Mitigation: Early user involvement, comprehensive training, continuous feedback loop.",
        "",
        "4. Risk: Integration Issues with Existing Systems",
           "   Mitigation: Detailed interface design, phased integration testing, fallback mechanisms.",
        "",
        "5. Risk: Security Incidents",
           "   Mitigation: Security design reviews, vulnerability assessments, incident response plan."
    ]
    apply_body_style(risk_box, risk_points, font_size=BODY_SIZE, para_spacing=Pt(10))
    # Adjust spacing for mitigation lines
    tf = risk_box.text_frame
    for i, p in enumerate(tf.paragraphs):
        if i > 1 and i % 3 == 0: # Mitigation lines
            p.level = 1
            p.space_before = Pt(0)
            p.space_after = Pt(10)
            if p.runs:
                 p.runs[0].font.size = Pt(14) # Slightly smaller for mitigation

    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="7/10")

def create_budget(prs):
    """予算計画のスライド (シンプル)"""
    slide_layout = prs.slide_layouts[5] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1),
        fill_color=ColorPalette.HEADING_BG
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(12), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Budget Plan & ROI"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT

    # 予算コンテンツ (2カラム風)
    # 左カラム: コスト
    cost_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5),
        Inches(5.5), Inches(5)
    )
    cost_tf = cost_box.text_frame
    cost_tf.clear() # Clear frame before adding content
    
    cost_p = cost_tf.paragraphs[0]
    cost_p.text = "Initial Investment"
    cost_run = cost_p.runs[0]
    cost_run.font.name = TITLE_FONT
    cost_run.font.size = SUBHEADING_SIZE
    cost_run.font.bold = True
    cost_p.space_after = Pt(12)

    initial_text = [
        "• Design & Development: ¥20M",
        "• Hardware & Cloud Setup: ¥5M",
        "• Data Migration & Testing: ¥6M",
        "• Training & Support: ¥4M",
        "• Total Initial Cost: ¥35M"
    ]
    # Add bullet points after the header
    for text in initial_text:
        p = cost_tf.add_paragraph()
        if text.startswith("• "):
            p.text = text[2:]
            p.level = 1
        else:
            p.text = text
            p.level = 0
        p.space_after = Pt(10)
        run = p.runs[0]
        run.font.name = BODY_FONT
        run.font.size = BODY_SIZE
        run.font.color.rgb = ColorPalette.TEXT

    run_p = cost_tf.add_paragraph()
    run_p.text = "Annual Running Costs"
    run_run = run_p.runs[0]
    run_run.font.name = TITLE_FONT
    run_run.font.size = SUBHEADING_SIZE
    run_run.font.bold = True
    run_p.space_before = Pt(18)
    run_p.space_after = Pt(12)

    running_text = [
        "• Cloud Infrastructure: ¥3M",
        "• Licensing Fees: ¥2M",
        "• Maintenance & Support: ¥3M",
        "• Total Annual Cost: ¥8M"
    ]
    # Add bullet points after the header
    for text in running_text:
        p = cost_tf.add_paragraph()
        if text.startswith("• "):
            p.text = text[2:]
            p.level = 1
        else:
            p.text = text
            p.level = 0
        p.space_after = Pt(10)
        run = p.runs[0]
        run.font.name = BODY_FONT
        run.font.size = BODY_SIZE
        run.font.color.rgb = ColorPalette.TEXT

    # 右カラム: ROI
    roi_box = slide.shapes.add_textbox(
        Inches(7), Inches(1.5),
        Inches(5.5), Inches(5)
    )
    roi_tf = roi_box.text_frame
    roi_tf.clear() # Clear frame before adding content

    roi_p = roi_tf.paragraphs[0]
    roi_p.text = "Return on Investment (ROI)"
    roi_run = roi_p.runs[0]
    roi_run.font.name = TITLE_FONT
    roi_run.font.size = SUBHEADING_SIZE
    roi_run.font.bold = True
    roi_p.space_after = Pt(12)

    roi_text = [
        "【Cost Savings】",
        "• Labor cost reduction (efficiency): ¥12M/year",
        "• System consolidation savings: ¥8M/year",
        "",
        "【Qualitative Benefits】",
        "• Faster decision-making",
        "• Improved customer satisfaction",
        "• Strategic advantage through data utilization",
        "",
        "【Payback Period】",
        "• Approx. 18 months"
    ]
    # Add text points after the header
    for i, text in enumerate(roi_text):
        p = roi_tf.add_paragraph()
        if text.startswith("【") and text.endswith("】"):
            p.text = text
            p.level = 0
            p.space_before = Pt(15) if i > 0 else Pt(0)
            run = p.runs[0]
            run.font.bold = True
            run.font.size = BODY_SIZE + Pt(2)
        elif text.startswith("• "):
            p.text = text[2:]
            p.level = 1
            run = p.runs[0]
            run.font.bold = False
        elif text.strip(): # Handle non-empty, non-special lines
            p.text = text
            p.level = 0
            run = p.runs[0]
            run.font.bold = False
        else: # Handle empty lines for spacing
            p.text = ""
            p.level = 0
            run = None # No run for empty lines

        p.space_after = Pt(10)
        if run:
            run.font.name = BODY_FONT
            run.font.size = BODY_SIZE
            run.font.color.rgb = ColorPalette.TEXT

    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="8/10")

def create_success_criteria(prs):
    """成功基準と評価方法のスライド (シンプル)"""
    slide_layout = prs.slide_layouts[5] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1),
        fill_color=ColorPalette.HEADING_BG
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(12), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "Success Criteria & Evaluation"
    header_p.alignment = PP_ALIGN.LEFT
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.HEADING_TEXT

    # 成功基準コンテンツ
    criteria_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5),
        Inches(11), Inches(5.5)
    )
    
    criteria_points = [
        "【Key Performance Indicators (KPIs)】",
        "",
        "System Performance Metrics:",
        "  • Response Time: < 2 seconds (peak)",
        "  • Availability: > 99.9%",
        "  • Concurrent Users: Support up to 300",
        "  • Backup Recovery Time: < 4 hours",
        "",
        "Business Impact Metrics:",
        "  • Process Time Reduction: 30%",
        "  • Customer Response Time Improvement: 50%",
        "  • Data Entry Error Reduction: 90%",
        "  • Paper Usage Reduction: 80%",
        "  • User Satisfaction: > 80%",
        "",
        "【Evaluation Method】",
        "• Quarterly performance measurement reports.",
        "• Monthly user satisfaction surveys.",
        "• Regular tracking of business efficiency metrics.",
        "• Continuous monitoring via real-time dashboards."
    ]
    apply_body_style(criteria_box, criteria_points, font_size=BODY_SIZE, para_spacing=Pt(10))
    tf = criteria_box.text_frame
    for i, p in enumerate(tf.paragraphs):
        # Make metric titles slightly bolder or distinct if needed
        if p.text.endswith(":") and not p.text.startswith("【"):
            if p.runs:
                p.runs[0].font.bold = True
                p.space_after = Pt(6)
        # Indent metric details
        if i in [3,4,5,6, 9,10,11,12,13]:
             p.level = 1
             p.space_before = Pt(0)
             p.space_after = Pt(6)
             if p.runs:
                 p.runs[0].font.size = Pt(14)

    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="9/10")

def create_conclusion(prs):
    """まとめと次のステップのスライド (ミニマル)"""
    slide_layout = prs.slide_layouts[5] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定 (黒背景)
    add_background(slide, prs, color=ColorPalette.HEADING_BG)
    
    # タイトル (白文字)
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5),
        Inches(11), Inches(1)
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "Conclusion & Next Steps"
    title_p.alignment = PP_ALIGN.LEFT
    title_run = title_p.runs[0]
    title_run.font.name = TITLE_FONT
    title_run.font.size = HEADING_SIZE
    title_run.font.bold = True
    title_run.font.color.rgb = ColorPalette.HEADING_TEXT

    # まとめコンテンツ (白文字)
    summary_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.8),
        Inches(11), Inches(1.5)
    )
    summary_text = [
        "【Proposal Summary】",
        "• Implement cloud-based system for 30% efficiency gain.",
        "• Phased 6-month rollout minimizes business disruption.",
        "• Investment: ¥35M initial, ¥8M annual. ROI within 18 months."
    ]
    apply_body_style(summary_box, summary_text, color=ColorPalette.HEADING_TEXT, para_spacing=Pt(12))

    # 次のステップコンテンツ (白文字)
    next_steps_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.5),
        Inches(11), Inches(1.5)
    )
    next_text = [
        "【Next Steps】",
        "• Final review and approval of proposal (within 1 week).",
        "• Project kick-off meeting (within 2 weeks of approval).",
        "• Commence detailed requirements definition (First week of April)."
    ]
    apply_body_style(next_steps_box, next_text, color=ColorPalette.HEADING_TEXT, para_spacing=Pt(12))

    # 連絡先 (白文字、シンプルに)
    contact_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.5),
        Inches(11), Inches(0.5)
    )
    contact_tf = contact_box.text_frame
    contact_p = contact_tf.paragraphs[0]
    contact_p.text = "Contact: Taro Yamada | yamada.taro@example.com | 03-1234-5678"
    contact_p.alignment = PP_ALIGN.LEFT
    contact_run = contact_p.runs[0]
    contact_run.font.name = BODY_FONT
    contact_run.font.size = BODY_SIZE
    contact_run.font.color.rgb = ColorPalette.FOOTER_BG # やや薄い白

    # フッターは不要 or ページ番号のみ
    # add_footer(slide, prs, text="", include_page_number=True, page_num="10/10")

if __name__ == "__main__":
    create_presentation()
