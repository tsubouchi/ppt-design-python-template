import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

# 洗練されたモダンなカラーパレットの定義
class ColorPalette:
    PRIMARY = RGBColor(45, 52, 54)         # ダークグレー（ほぼブラック）
    SECONDARY = RGBColor(85, 239, 196)     # ミントグリーン
    ACCENT1 = RGBColor(129, 236, 236)      # ターコイズ
    ACCENT2 = RGBColor(250, 177, 160)      # サーモンピンク
    ACCENT3 = RGBColor(116, 185, 255)      # ソフトブルー
    DARK = RGBColor(30, 39, 46)            # ディープダークグレー
    LIGHT = RGBColor(245, 246, 250)        # オフホワイト
    GRAY = RGBColor(206, 214, 224)         # ライトグレー
    GRADIENT_START = RGBColor(85, 239, 196) # グラデーション開始色
    GRADIENT_END = RGBColor(129, 236, 236)  # グラデーション終了色

# 洗練されたフォント設定
TITLE_FONT = 'Avenir Next'  # よりモダンなフォント
BODY_FONT = 'Avenir'
TITLE_SIZE = Pt(40)         # より大きなタイトル
SUBTITLE_SIZE = Pt(24)
HEADING_SIZE = Pt(32)
SUBHEADING_SIZE = Pt(22)
BODY_SIZE = Pt(18)

def create_presentation():
    prs = Presentation()
    
    # スライドサイズを16:9に設定
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # スライドを作成
    create_title_slide(prs)
    create_executive_summary(prs)
    create_current_analysis(prs)
    create_proposal(prs)
    create_schedule(prs)
    create_team_structure(prs)
    create_risk_management(prs)
    create_budget(prs)
    create_success_criteria(prs)
    create_conclusion(prs)
    
    # プレゼンテーションを保存
    prs.save('IT_Project_Proposal.pptx')
    print("プレゼンテーションが作成されました: IT_Project_Proposal.pptx")

def apply_title_style(title_shape, text, font_size=TITLE_SIZE, color=ColorPalette.DARK):
    """タイトルのスタイルを適用する"""
    title_shape.text = text
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_run = title_para.runs[0]
    title_run.font.name = TITLE_FONT
    title_run.font.size = font_size
    title_run.font.bold = True
    title_run.font.color.rgb = color

def apply_body_style(body_shape, text_list, font_size=BODY_SIZE, color=ColorPalette.DARK):
    """本文のスタイルを適用する"""
    tf = body_shape.text_frame
    tf.clear()
    
    for i, text in enumerate(text_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = text
        p.level = 0 if not text.startswith('•') else 1
        p.space_after = Pt(12)
        
        # テキストを設定した後、新しいrunを追加して書式設定
        if not p.runs:
            run = p.add_run()
            run.text = text
        else:
            run = p.runs[0]
            
        run.font.name = BODY_FONT
        run.font.size = font_size
        run.font.color.rgb = color

def add_decorative_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None, line_width=Pt(1), shadow=False, transparency=0):
    """洗練された装飾的な図形を追加する"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    
    # 塗りつぶし設定
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.fill.transparency = transparency
    
    # 線の設定
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()  # 線なし
    
    # 影の設定
    if shadow:
        shape.shadow.inherit = False
        shape.shadow.visible = True
        shape.shadow.blur_radius = Pt(5)
        shape.shadow.distance = Pt(3)
        shape.shadow.angle = 45
        # python-pptxの最新バージョンでは影の色の設定方法が異なる場合があります
        try:
            shape.shadow.color.rgb = RGBColor(0, 0, 0)
            shape.shadow.transparency = 0.6
        except AttributeError:
            # 代替方法：影の設定をスキップ
            pass
    
    return shape

def add_gradient_shape(slide, shape_type, left, top, width, height, start_color, end_color, direction='horizontal', line_color=None):
    """グラデーション効果のある図形を追加する"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    
    # グラデーション設定
    shape.fill.gradient()
    shape.fill.gradient_stops[0].position = 0
    shape.fill.gradient_stops[0].color.rgb = start_color
    shape.fill.gradient_stops[1].position = 1
    shape.fill.gradient_stops[1].color.rgb = end_color
    
    # グラデーション方向
    if direction == 'horizontal':
        shape.fill.gradient_angle = 90
    else:  # vertical
        shape.fill.gradient_angle = 0
    
    # 線の設定
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()  # 線なし
    
    return shape

def create_title_slide(prs):
    """洗練された表紙スライドの作成"""
    slide_layout = prs.slide_layouts[0]  # タイトルスライド
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景全体にグラデーションを適用
    background = add_gradient_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height,
        ColorPalette.LIGHT, ColorPalette.LIGHT,
        'vertical'
    )
    
    # 左側に装飾的なグラデーションバー
    left_bar = add_gradient_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(2.5), prs.slide_height,
        ColorPalette.GRADIENT_START, ColorPalette.GRADIENT_END,
        'vertical'
    )
    
    # 右下に装飾的な要素
    bottom_shape = add_decorative_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8), Inches(6),
        Inches(5.33), Inches(1.5),
        ColorPalette.DARK, None, Pt(0), True, 0.1
    )
    
    # 装飾的な円形要素を追加
    circle1 = add_decorative_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(11), Inches(1),
        Inches(1.5), Inches(1.5),
        ColorPalette.ACCENT1, None, Pt(0), False, 0.5
    )
    
    circle2 = add_decorative_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(10.5), Inches(2),
        Inches(0.8), Inches(0.8),
        ColorPalette.ACCENT2, None, Pt(0), False, 0.6
    )
    
    # タイトル用のカスタムテキストボックス
    title_box = slide.shapes.add_textbox(
        Inches(3), Inches(2.5),
        Inches(9), Inches(1.5)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.LEFT
    title_p.text = "IT開発・システム導入"
    title_run = title_p.runs[0]
    title_run.font.name = TITLE_FONT
    title_run.font.size = TITLE_SIZE
    title_run.font.bold = True
    title_run.font.color.rgb = ColorPalette.DARK
    
    subtitle_p = title_tf.add_paragraph()
    subtitle_p.text = "プロジェクト計画書"
    subtitle_run = subtitle_p.runs[0]
    subtitle_run.font.name = TITLE_FONT
    subtitle_run.font.size = SUBTITLE_SIZE
    subtitle_run.font.bold = True
    subtitle_run.font.color.rgb = ColorPalette.DARK
    
    # 日付
    date_box = slide.shapes.add_textbox(
        Inches(3), Inches(4.5),
        Inches(4), Inches(0.5)
    )
    date_tf = date_box.text_frame
    date_p = date_tf.paragraphs[0]
    date_p.text = "2025年3月30日"
    date_run = date_p.runs[0]
    date_run.font.name = BODY_FONT
    date_run.font.size = Pt(16)
    date_run.font.color.rgb = ColorPalette.DARK
    
    # 会社名
    company_box = slide.shapes.add_textbox(
        Inches(3), Inches(5),
        Inches(4), Inches(0.5)
    )
    company_tf = company_box.text_frame
    company_p = company_tf.paragraphs[0]
    company_p.text = "株式会社〇〇〇〇"
    company_run = company_p.runs[0]
    company_run.font.name = BODY_FONT
    company_run.font.size = Pt(16)
    company_run.font.bold = True
    company_run.font.color.rgb = ColorPalette.DARK

def create_executive_summary(prs):
    """洗練されたエグゼクティブサマリーのスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景全体に薄いグラデーションを適用
    background = add_gradient_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height,
        ColorPalette.LIGHT, ColorPalette.LIGHT,
        'vertical'
    )
    
    # 上部にグラデーションのヘッダーを追加
    header = add_gradient_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2),
        ColorPalette.DARK, ColorPalette.PRIMARY,
        'horizontal'
    )
    
    # 装飾的な要素を追加
    accent_shape1 = add_decorative_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(1.2),
        Inches(0.2), Inches(6.3),
        ColorPalette.ACCENT1, None, Pt(0), False, 0.2
    )
    
    accent_shape2 = add_decorative_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(12), Inches(1.5),
        Inches(0.8), Inches(0.8),
        ColorPalette.ACCENT3, None, Pt(0), False, 0.5
    )
    
    # タイトル用のカスタムテキストボックス
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12), Inches(0.6)
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.LEFT
    title_p.text = "エグゼクティブサマリー"
    title_run = title_p.runs[0]
    title_run.font.name = TITLE_FONT
    title_run.font.size = HEADING_SIZE
    title_run.font.bold = True
    title_run.font.color.rgb = ColorPalette.LIGHT
    
    # 内容用のカスタムテキストボックス
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.8),
        Inches(11), Inches(5)
    )
    
    summary_points = [
        "• プロジェクトの目的: 現行の業務システムを刷新し、業務効率を30%向上",
        "• 主要な提案内容: クラウドベースの統合管理システムの導入",
        "• 期待される効果: 年間コスト削減2,000万円、顧客対応時間50%短縮",
        "• 実施期間: 2025年4月〜2025年9月（6ヶ月間）",
        "• 予算概要: 初期投資3,500万円、年間運用コスト800万円",
        "• ROI: 導入後18ヶ月で投資回収見込み"
    ]
    apply_body_style(content_box, summary_points)
    
    # 装飾的な図形を右下に追加
    corner_shape = add_decorative_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(11), Inches(6),
        Inches(2.33), Inches(1.5),
        ColorPalette.ACCENT2, None, Pt(0), True, 0.8
    )

def create_current_analysis(prs):
    """洗練された現状分析と課題のスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景全体に薄いグラデーションを適用
    background = add_gradient_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height,
        ColorPalette.LIGHT, ColorPalette.LIGHT,
        'vertical'
    )
    
    # 上部にグラデーションのヘッダーを追加
    header = add_gradient_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2),
        ColorPalette.DARK, ColorPalette.PRIMARY,
        'horizontal'
    )
    
    # タイトル用のカスタムテキストボックス
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12), Inches(0.6)
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.LEFT
    title_p.text = "現状分析と課題"
    title_run = title_p.runs[0]
    title_run.font.name = TITLE_FONT
    title_run.font.size = HEADING_SIZE
    title_run.font.bold = True
    title_run.font.color.rgb = ColorPalette.LIGHT
    
    # 左側のパネル: 現状
    left_panel = add_gradient_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5),
        Inches(6), Inches(5.5),
        ColorPalette.LIGHT, ColorPalette.LIGHT,
        'vertical'
    )
    
    # 左側パネルのヘッダー
    left_header = add_gradient_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5),
        Inches(6), Inches(0.6),
        ColorPalette.ACCENT1, ColorPalette.ACCENT1,
        'horizontal'
    )
    
    # 左側ヘッダーテキスト
    left_title = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.6),
        Inches(5.6), Inches(0.4)
    )
    left_title_tf = left_title.text_frame
    left_title_p = left_title_tf.paragraphs[0]
    left_title_p.text = "現在のシステム状況"
    left_title_p.alignment = PP_ALIGN.CENTER
    left_title_run = left_title_p.runs[0]
    left_title_run.font.name = TITLE_FONT
    left_title_run.font.size = SUBHEADING_SIZE
    left_title_run.font.bold = True
    left_title_run.font.color.rgb = ColorPalette.LIGHT
    
    # 左側: 現状
    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(5.6), Inches(4.5))
    current_state = [
        "• 導入から8年経過した基幹システム",
        "• 複数のシステムが連携せず、二重入力が発生",
        "• レガシーシステムによるメンテナンスコスト増加",
        "• オンプレミス環境でのリソース制約",
        "• モバイル対応していないため外出先での業務に制約"
    ]
    apply_body_style(left_box, current_state)
    
    # 右側のパネル: 課題
    right_panel = add_gradient_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.5),
        Inches(6), Inches(5.5),
        ColorPalette.LIGHT, ColorPalette.LIGHT,
        'vertical'
    )
    
    # 右側パネルのヘッダー
    right_header = add_gradient_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.5),
        Inches(6), Inches(0.6),
        ColorPalette.ACCENT2, ColorPalette.ACCENT2,
        'horizontal'
    )
    
    # 右側ヘッダーテキスト
    right_title = slide.shapes.add_textbox(
        Inches(7), Inches(1.6),
        Inches(5.6), Inches(0.4)
    )
    right_title_tf = right_title.text_frame
    right_title_p = right_title_tf.paragraphs[0]
    right_title_p.text = "解決すべき課題"
    right_title_p.alignment = PP_ALIGN.CENTER
    right_title_run = right_title_p.runs[0]
    right_title_run.font.name = TITLE_FONT
    right_title_run.font.size = SUBHEADING_SIZE
    right_title_run.font.bold = True
    right_title_run.font.color.rgb = ColorPalette.LIGHT
    
    # 右側: 課題
    right_box = slide.shapes.add_textbox(Inches(7), Inches(2.3), Inches(5.6), Inches(4.5))
    challenges = [
        "• データの一元管理と業務プロセスの標準化",
        "• システム間連携の自動化による二重作業の排除",
        "• クラウド環境への移行によるコスト最適化",
        "• モバイル対応によるリモートワーク環境の整備",
        "• セキュリティ強化とコンプライアンス対応"
    ]
    apply_body_style(right_box, challenges)
    
    # 装飾的な要素を追加
    accent_shape = add_decorative_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(6.4), Inches(4),
        Inches(0.8), Inches(0.8),
        ColorPalette.ACCENT3, None, Pt(0), True, 0.7
    )

def create_proposal(prs):
    """提案内容のスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 装飾的な要素を追加
    top_shape = add_decorative_shape(
        slide, MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        prs.slide_width, Inches(0.5),
        ColorPalette.PRIMARY
    )
    
    # タイトル
    title_shape = slide.shapes.title
    apply_title_style(title_shape, "提案内容: クラウド統合管理システム", HEADING_SIZE)
    
    # 提案内容の説明
    content = slide.placeholders[1]
    proposal_points = [
        "【システム概要】",
        "• クラウドベースの統合業務管理プラットフォーム",
        "• マイクロサービスアーキテクチャによる柔軟な拡張性",
        "• APIによる外部システムとのシームレスな連携",
        "• レスポンシブデザインによるマルチデバイス対応",
        "• AIを活用した業務自動化と予測分析機能",
        "",
        "【主要機能】",
        "• 顧客情報管理・案件管理・進捗管理の統合",
        "• リアルタイムダッシュボードとレポーティング",
        "• ワークフロー自動化とタスク管理",
        "• 権限管理とセキュリティ制御"
    ]
    apply_body_style(content, proposal_points)
    
    # 装飾的な図形を右下に追加
    corner_shape = add_decorative_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(12), Inches(6.5), 
        Inches(1.33), Inches(1),
        ColorPalette.ACCENT3
    )

def create_schedule(prs):
    """導入スケジュールのスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 装飾的な要素を追加
    top_shape = add_decorative_shape(
        slide, MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        prs.slide_width, Inches(0.5),
        ColorPalette.PRIMARY
    )
    
    # タイトル
    title_shape = slide.shapes.title
    apply_title_style(title_shape, "導入スケジュール（6ヶ月計画）", HEADING_SIZE)
    
    # スケジュール表の代わりにテキストで表現
    content = slide.placeholders[1]
    schedule_points = [
        "【フェーズ1: 要件定義・設計（4月〜5月）】",
        "• 業務要件の詳細ヒアリングと分析",
        "• システム設計とアーキテクチャ確定",
        "• データ移行計画の策定",
        "",
        "【フェーズ2: 開発・構築（5月〜7月）】",
        "• システム基盤構築とコア機能の開発",
        "• 外部システム連携の実装",
        "• ユーザーインターフェース開発",
        "",
        "【フェーズ3: テスト・移行（7月〜8月）】",
        "• 単体・結合テストの実施",
        "• ユーザー受け入れテスト",
        "• データ移行とシステム切り替え準備",
        "",
        "【フェーズ4: 本番稼働・安定化（9月）】",
        "• 段階的な本番リリース",
        "• ユーザートレーニングの実施",
        "• 運用体制の確立とサポート"
    ]
    apply_body_style(content, schedule_points, Pt(16))  # フォントサイズを少し小さく

def create_team_structure(prs):
    """実施体制のスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 装飾的な要素を追加
    top_shape = add_decorative_shape(
        slide, MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        prs.slide_width, Inches(0.5),
        ColorPalette.PRIMARY
    )
    
    # タイトル
    title_shape = slide.shapes.title
    apply_title_style(title_shape, "実施体制", HEADING_SIZE)
    
    # 実施体制の説明
    content = slide.placeholders[1]
    team_points = [
        "【プロジェクト推進体制】",
        "• プロジェクトスポンサー: 経営企画部長",
        "• プロジェクトマネージャー: IT部門 課長",
        "• テクニカルリード: システム開発チーム リーダー",
        "• 業務プロセス担当: 各部門代表者",
        "",
        "【役割と責任】",
        "• 要件定義・設計: 弊社コンサルタント + お客様業務担当者",
        "• システム開発: 弊社エンジニアチーム（5名）",
        "• テスト・品質保証: 弊社QAチーム + お客様検証担当者",
        "• 導入・トレーニング: 弊社導入支援チーム",
        "",
        "【コミュニケーション体制】",
        "• 週次進捗会議（オンライン）",
        "• 月次ステアリングコミッティ（対面）",
        "• 日次スクラムミーティング（開発チーム）"
    ]
    apply_body_style(content, team_points)
    
    # 装飾的な図形を左下に追加
    corner_shape = add_decorative_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(0), Inches(6.5), 
        Inches(1.33), Inches(1),
        ColorPalette.ACCENT1
    )

def create_risk_management(prs):
    """リスク管理計画のスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 装飾的な要素を追加
    top_shape = add_decorative_shape(
        slide, MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        prs.slide_width, Inches(0.5),
        ColorPalette.PRIMARY
    )
    
    # タイトル
    title_shape = slide.shapes.title
    apply_title_style(title_shape, "リスク管理計画", HEADING_SIZE)
    
    # リスク管理の説明
    content = slide.placeholders[1]
    risk_points = [
        "【主要リスクと対策】",
        "",
        "• リスク1: 要件定義の不足・変更による開発遅延",
        "  対策: アジャイル開発手法の採用、定期的な要件レビュー",
        "",
        "• リスク2: データ移行時のデータ欠損・不整合",
        "  対策: 事前データクレンジング、段階的移行、二重検証体制",
        "",
        "• リスク3: ユーザー受け入れの低さ",
        "  対策: 早期からのユーザー参加、充実したトレーニング計画",
        "",
        "• リスク4: 既存システムとの連携不具合",
        "  対策: 詳細なインターフェース設計、段階的な連携テスト",
        "",
        "• リスク5: セキュリティインシデント",
        "  対策: セキュリティ設計レビュー、脆弱性診断の実施"
    ]
    apply_body_style(content, risk_points)
    
    # 装飾的な図形を右下に追加
    corner_shape = add_decorative_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(12), Inches(6.5), 
        Inches(1.33), Inches(1),
        ColorPalette.ACCENT2
    )

def create_budget(prs):
    """予算計画のスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 装飾的な要素を追加
    top_shape = add_decorative_shape(
        slide, MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        prs.slide_width, Inches(0.5),
        ColorPalette.PRIMARY
    )
    
    # タイトル
    title_shape = slide.shapes.title
    apply_title_style(title_shape, "予算計画", HEADING_SIZE)
    
    # 予算の説明
    content = slide.placeholders[1]
    budget_points = [
        "【初期導入コスト】",
        "• システム設計・開発費: 2,000万円",
        "• ハードウェア・クラウド環境構築: 500万円",
        "• データ移行・テスト: 600万円",
        "• トレーニング・導入支援: 400万円",
        "• 初期費用合計: 3,500万円",
        "",
        "【ランニングコスト（年間）】",
        "• クラウドインフラ利用料: 300万円",
        "• ライセンス費用: 200万円",
        "• 保守・サポート費: 300万円",
        "• 年間運用コスト合計: 800万円",
        "",
        "【投資対効果（ROI）】",
        "• 業務効率化による人件費削減: 年間1,200万円",
        "• システム統合によるコスト削減: 年間800万円",
        "• 投資回収期間: 約18ヶ月"
    ]
    apply_body_style(content, budget_points)
    
    # 装飾的な図形を左下に追加
    corner_shape = add_decorative_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(0), Inches(6.5), 
        Inches(1.33), Inches(1),
        ColorPalette.ACCENT3
    )

def create_success_criteria(prs):
    """成功基準と評価方法のスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 装飾的な要素を追加
    top_shape = add_decorative_shape(
        slide, MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        prs.slide_width, Inches(0.5),
        ColorPalette.PRIMARY
    )
    
    # タイトル
    title_shape = slide.shapes.title
    apply_title_style(title_shape, "成功基準と評価方法", HEADING_SIZE)
    
    # 成功基準の説明
    content = slide.placeholders[1]
    criteria_points = [
        "【KPI（主要業績評価指標）】",
        "",
        "• システムパフォーマンス指標",
        "  - システム応答時間: 2秒以内（ピーク時）",
        "  - システム可用性: 99.9%以上",
        "  - 同時接続ユーザー: 最大300名をサポート",
        "",
        "• ビジネス効果指標",
        "  - 業務処理時間: 30%削減",
        "  - 顧客対応時間: 50%短縮",
        "  - データ入力エラー: 90%削減",
        "  - ペーパーレス化: 紙使用量80%削減",
        "",
        "• 測定方法",
        "  - 四半期ごとのパフォーマンス測定",
        "  - 月次ユーザー満足度調査",
        "  - 業務効率化指標の定期測定"
    ]
    apply_body_style(content, criteria_points)
    
    # 装飾的な図形を右下に追加
    corner_shape = add_decorative_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(12), Inches(6.5), 
        Inches(1.33), Inches(1),
        ColorPalette.ACCENT1
    )

def create_conclusion(prs):
    """まとめと次のステップのスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 装飾的な要素を追加
    top_shape = add_decorative_shape(
        slide, MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        prs.slide_width, Inches(0.5),
        ColorPalette.PRIMARY
    )
    
    # 背景に装飾的な図形を追加
    bg_shape = add_decorative_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(0), Inches(6.5), 
        prs.slide_width, Inches(1),
        ColorPalette.PRIMARY
    )
    
    # タイトル
    title_shape = slide.shapes.title
    apply_title_style(title_shape, "まとめと次のステップ", HEADING_SIZE)
    
    # まとめの説明
    content = slide.placeholders[1]
    conclusion_points = [
        "【提案のまとめ】",
        "• クラウドベースの統合管理システム導入により業務効率を30%向上",
        "• 6ヶ月間の段階的な導入計画で業務への影響を最小化",
        "• 初期投資3,500万円、年間運用コスト800万円、18ヶ月でROI達成",
        "",
        "【次のステップ】",
        "• 提案内容の最終確認と承認（1週間以内）",
        "• キックオフミーティングの開催（承認後2週間以内）",
        "• 詳細要件定義の開始（4月第1週）",
        "",
        "【連絡先】",
        "• プロジェクト担当: 山田太郎",
        "• メール: yamada.taro@example.com",
        "• 電話: 03-1234-5678"
    ]
    apply_body_style(content, conclusion_points)
    
    # 装飾的な図形を左下に追加
    corner_shape = add_decorative_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(0), Inches(5), 
        Inches(0.5), Inches(1.5),
        ColorPalette.ACCENT1
    )

if __name__ == "__main__":
    create_presentation()