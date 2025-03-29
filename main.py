import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL

# より洗練されたモダンなカラーパレットの定義
class ColorPalette:
    # メインカラー
    PRIMARY = RGBColor(28, 30, 39)         # ディープネイビー（アクセントより暗い）
    SECONDARY = RGBColor(243, 243, 243)    # オフホワイト（背景用）
    
    # アクセントカラー
    ACCENT1 = RGBColor(75, 192, 192)       # ティール
    ACCENT2 = RGBColor(255, 159, 64)       # オレンジ
    ACCENT3 = RGBColor(153, 102, 255)      # パープル
    ACCENT4 = RGBColor(54, 162, 235)       # ブルー
    
    # サポートカラー
    DARK = RGBColor(22, 22, 29)            # ほぼブラック
    LIGHT = RGBColor(250, 250, 250)        # ホワイト
    GRAY = RGBColor(130, 138, 153)         # ミディアムグレー
    LIGHT_GRAY = RGBColor(240, 242, 245)   # 薄いグレー

# モダンでエレガントなフォント設定
TITLE_FONT = 'Montserrat'      # モダンなサンセリフ
BODY_FONT = 'Noto Sans'        # 読みやすいサンセリフ（日本語対応）
ALT_FONT = 'Helvetica Neue'    # 代替フォント

# フォントサイズ定義
TITLE_SIZE = Pt(44)            # 大きくはっきりとしたタイトル
HEADING_SIZE = Pt(36)          # セクションタイトル
SUBHEADING_SIZE = Pt(24)       # サブセクションタイトル
BODY_SIZE = Pt(18)             # 本文
CAPTION_SIZE = Pt(14)          # キャプションやフッター用

def create_presentation():
    prs = Presentation()
    
    # スライドサイズを16:9に設定
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # スライドを作成
    create_title_slide(prs)
    create_executive_summary(prs)
    create_current_analysis(prs)
    create_proposal(prs)
    create_schedule(prs)
    create_team_structure(prs)
    create_risk_management(prs)
    create_budget(prs)
    create_success_criteria(prs)
    create_conclusion(prs)
    
    # プレゼンテーションを保存
    prs.save('sample1.pptx')
    print("洗練されたプレゼンテーションが作成されました: sample1.pptx")

def apply_title_style(title_shape, text, font_size=TITLE_SIZE, color=ColorPalette.PRIMARY, align=PP_ALIGN.LEFT, bold=True):
    """タイトルのスタイルを適用する"""
    title_shape.text = text
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.alignment = align
    title_run = title_para.runs[0]
    title_run.font.name = TITLE_FONT
    title_run.font.size = font_size
    title_run.font.bold = bold
    title_run.font.color.rgb = color

def apply_body_style(body_shape, text_list, font_size=BODY_SIZE, color=ColorPalette.DARK, para_spacing=Pt(10)):
    """本文のスタイルを適用する"""
    tf = body_shape.text_frame
    tf.clear()

    # 最初の段落を確保
    if not tf.paragraphs:
        tf.add_paragraph()
    p = tf.paragraphs[0]

    for i, original_text in enumerate(text_list):
        text_to_set = original_text  # スタイル適用用のテキスト

        # 2番目以降の要素のために新しい段落を追加
        if i > 0:
            p = tf.add_paragraph()

        # 空文字列または空白のみの文字列を処理
        if not original_text.strip():
            p.text = ""  # 空のテキストを設定
            p.space_after = para_spacing
            continue  # 次の要素へ

        # レベル決定とテキスト変更（割り当て前）
        level = 0
        space_before = Pt(0)
        is_heading = False
        is_bullet1 = False
        is_bullet2 = False

        if original_text.startswith('【') and original_text.endswith('】'):
            level = 0
            space_before = Pt(15) if i > 0 else Pt(0)
            is_heading = True
        elif original_text.startswith('• '):
            level = 1
            text_to_set = original_text[2:]
            is_bullet1 = True
        elif original_text.startswith('  - '):
            level = 2
            text_to_set = original_text[4:]
            is_bullet2 = True
        else:
            level = 0

        # テキスト、レベル、スペースを設定
        p.text = text_to_set
        p.level = level
        p.space_after = para_spacing
        p.space_before = space_before

        # runが存在するか確認してスタイルを適用
        if p.runs:
            run = p.runs[0]
            run.font.name = BODY_FONT
            run.font.size = font_size

            # 強調表示などの特定のスタイル
            if is_heading:
                run.font.bold = True
                run.font.size = font_size + Pt(2)
                run.font.color.rgb = ColorPalette.ACCENT1
            else:
                run.font.bold = False
                run.font.color.rgb = color
        # else:
            # テキストが空でないのにrunが作成されない場合（通常は発生しない）
            # print(f"警告: テキスト '{text_to_set}' のrunが作成されませんでした")
            # pass

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, 
              line_width=Pt(1), shadow=False, transparency=0, gradient_to=None, text=None):
    """洗練された図形を追加する"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    
    # 塗りつぶし設定
    if fill_color:
        if gradient_to:
            # グラデーション
            shape.fill.gradient()
            shape.fill.gradient_stops[0].position = 0
            shape.fill.gradient_stops[0].color.rgb = fill_color
            shape.fill.gradient_stops[1].position = 1
            shape.fill.gradient_stops[1].color.rgb = gradient_to
            shape.fill.gradient_angle = 90  # 左から右へのグラデーション
        else:
            # 単色
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            shape.fill.transparency = transparency
    else:
        # 塗りつぶしなし
        shape.fill.background()
    
    # 線の設定
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()  # 線なし
    
    # 影の設定
    if shadow:
        shape.shadow.inherit = False
        shape.shadow.visible = True
        shape.shadow.blur_radius = Pt(5)
        shape.shadow.distance = Pt(3)
        shape.shadow.angle = 45
        try:
            shape.shadow.color.rgb = RGBColor(0, 0, 0)
            shape.shadow.transparency = 0.7
        except AttributeError:
            pass
    
    # テキストがある場合は追加
    if text:
        shape.text = text
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = TITLE_FONT
        run.font.size = SUBHEADING_SIZE
        run.font.color.rgb = ColorPalette.LIGHT
    
    return shape

def add_background(slide, prs, type="solid", color=ColorPalette.SECONDARY, gradient_to=None):
    """スライドの背景を設定する"""
    if type == "solid":
        # 単色背景
        background = add_shape(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height,
            fill_color=color
        )
    elif type == "gradient":
        # グラデーション背景
        background = add_shape(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height,
            fill_color=color, gradient_to=gradient_to
        )
    elif type == "accent_bar":
        # アクセントバー付き背景
        background = add_shape(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height,
            fill_color=color
        )
        # 左サイドにアクセントバー
        accent_bar = add_shape(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.5), prs.slide_height,
            fill_color=gradient_to
        )
    
    return background

def add_footer(slide, prs, text="株式会社〇〇〇〇 | IT開発プロジェクト計画書 | 2025年3月", include_page_number=True, page_num=None):
    """フッターを追加する"""
    footer_shape = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), prs.slide_height - Inches(0.4),
        prs.slide_width, Inches(0.4),
        fill_color=ColorPalette.LIGHT_GRAY, transparency=0.5
    )
    
    # フッターテキスト
    footer_text = slide.shapes.add_textbox(
        Inches(0.5), prs.slide_height - Inches(0.35),
        prs.slide_width - Inches(1.5), Inches(0.3)
    )
    tf = footer_text.text_frame
    p = tf.paragraphs[0]
    
    if include_page_number and page_num:
        p.text = f"{text} | {page_num}"
    else:
        p.text = text
    
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = BODY_FONT
    run.font.size = CAPTION_SIZE
    run.font.color.rgb = ColorPalette.GRAY
    
    return footer_shape

def create_title_slide(prs):
    """洗練された表紙スライドの作成"""
    slide_layout = prs.slide_layouts[0]  # タイトルスライド
    slide = prs.slides.add_slide(slide_layout)
    
    # モダンな全画面グラデーション背景
    background = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height,
        fill_color=ColorPalette.PRIMARY, gradient_to=ColorPalette.DARK
    )
    
    # 装飾的な要素 - タイトル上部の線
    accent_line = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(2), Inches(2.2),
        Inches(2), Inches(0.1),
        fill_color=ColorPalette.ACCENT1
    )
    
    # 装飾的な図形 - 抽象的な要素
    decorative_circle1 = add_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(10), Inches(1),
        Inches(2), Inches(2),
        fill_color=ColorPalette.ACCENT1, transparency=0.85
    )
    
    decorative_circle2 = add_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(11), Inches(1.5),
        Inches(3), Inches(3),
        fill_color=ColorPalette.ACCENT3, transparency=0.9
    )
    
    decorative_circle3 = add_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(9.5), Inches(2),
        Inches(1.5), Inches(1.5),
        fill_color=ColorPalette.ACCENT2, transparency=0.8
    )
    
    # タイトル用のテキストボックス
    title_box = slide.shapes.add_textbox(
        Inches(2), Inches(2.5),
        Inches(9), Inches(2)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    
    # メインタイトル
    title_p = title_tf.paragraphs[0]
    title_p.text = "IT開発・システム導入"
    title_p.alignment = PP_ALIGN.LEFT
    title_run = title_p.runs[0]
    title_run.font.name = TITLE_FONT
    title_run.font.size = TITLE_SIZE
    title_run.font.bold = True
    title_run.font.color.rgb = ColorPalette.LIGHT
    
    # サブタイトル
    subtitle_p = title_tf.add_paragraph()
    subtitle_p.text = "プロジェクト計画書"
    subtitle_p.alignment = PP_ALIGN.LEFT
    subtitle_p.space_before = Pt(10)
    subtitle_run = subtitle_p.runs[0]
    subtitle_run.font.name = TITLE_FONT
    subtitle_run.font.size = SUBHEADING_SIZE
    subtitle_run.font.bold = True
    subtitle_run.font.color.rgb = ColorPalette.LIGHT
    
    # 日付
    date_box = slide.shapes.add_textbox(
        Inches(2), Inches(5),
        Inches(9), Inches(0.5)
    )
    date_tf = date_box.text_frame
    date_p = date_tf.paragraphs[0]
    date_p.text = "2025年3月30日"
    date_p.alignment = PP_ALIGN.LEFT
    date_run = date_p.runs[0]
    date_run.font.name = BODY_FONT
    date_run.font.size = BODY_SIZE
    date_run.font.color.rgb = ColorPalette.GRAY
    
    # 会社名
    company_box = slide.shapes.add_textbox(
        Inches(2), Inches(5.5),
        Inches(9), Inches(0.5)
    )
    company_tf = company_box.text_frame
    company_p = company_tf.paragraphs[0]
    company_p.text = "株式会社〇〇〇〇"
    company_p.alignment = PP_ALIGN.LEFT
    company_run = company_p.runs[0]
    company_run.font.name = BODY_FONT
    company_run.font.size = BODY_SIZE
    company_run.font.bold = True
    company_run.font.color.rgb = ColorPalette.LIGHT

def create_executive_summary(prs):
    """洗練されたエグゼクティブサマリーのスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs, "solid", ColorPalette.SECONDARY)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2),
        fill_color=ColorPalette.PRIMARY
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(1), Inches(0.3),
        Inches(11), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "エグゼクティブサマリー"
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.LIGHT
    
    # 左側のアクセントライン
    accent_line = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(0.1), Inches(5),
        fill_color=ColorPalette.ACCENT1
    )
    
    # 装飾的な要素 - 抽象的な図形
    decorative_shape1 = add_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(11), Inches(5.5),
        Inches(2), Inches(2),
        fill_color=ColorPalette.ACCENT1, transparency=0.9
    )
    
    decorative_shape2 = add_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(10.5), Inches(6),
        Inches(1), Inches(1),
        fill_color=ColorPalette.ACCENT3, transparency=0.8
    )
    
    # コンテンツパネル
    content_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(1.8),
        Inches(10), Inches(4.5),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.LIGHT_GRAY, line_width=Pt(0.75),
        shadow=True
    )
    
    # コンテンツ
    content_box = slide.shapes.add_textbox(
        Inches(2), Inches(2),
        Inches(9), Inches(4)
    )
    
    summary_points = [
        "【プロジェクトの目的】",
        "• 現行の業務システムを刷新し、業務効率を30%向上",
        "• デジタルトランスフォーメーションの基盤構築",
        "",
        "【主要な提案内容】",
        "• クラウドベースの統合管理システムの導入",
        "• AI活用による業務自動化と予測分析の実現",
        "",
        "【期待される効果】",
        "• 年間コスト削減2,000万円、顧客対応時間50%短縮",
        "• データ駆動型意思決定の実現とビジネス機会の拡大"
    ]
    apply_body_style(content_box, summary_points)
    
    # 予算・期間の情報パネル
    info_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.5), Inches(4.5),
        Inches(5), Inches(1.5),
        fill_color=ColorPalette.ACCENT4, transparency=0.1,
        line_color=ColorPalette.ACCENT4, line_width=Pt(0.75)
    )
    
    # 情報パネルのコンテンツ
    info_box = slide.shapes.add_textbox(
        Inches(6.8), Inches(4.7),
        Inches(4.5), Inches(1.2)
    )
    
    info_points = [
        "• 実施期間: 2025年4月〜2025年9月（6ヶ月間）",
        "• 予算概要: 初期投資3,500万円、年間運用コスト800万円",
        "• ROI: 導入後18ヶ月で投資回収見込み"
    ]
    apply_body_style(info_box, info_points, BODY_SIZE, ColorPalette.DARK, Pt(8))
    
    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="2/10")

def create_current_analysis(prs):
    """洗練された現状分析と課題のスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs, "solid", ColorPalette.SECONDARY)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2),
        fill_color=ColorPalette.PRIMARY
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(1), Inches(0.3),
        Inches(11), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "現状分析と課題"
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.LIGHT
    
    # 左側: 現状パネル
    current_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(5.5), Inches(5),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.LIGHT_GRAY, line_width=Pt(0.75),
        shadow=True
    )
    
    # 現状パネルのヘッダー
    current_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(5.5), Inches(0.7),
        fill_color=ColorPalette.ACCENT1,
        text="現在のシステム状況"
    )
    
    # 現状の内容
    current_box = slide.shapes.add_textbox(
        Inches(1.3), Inches(2.4),
        Inches(5), Inches(4)
    )
    
    current_state = [
        "• 導入から8年経過した基幹システム",
        "• 複数のシステムが連携せず、二重入力が発生",
        "• レガシーシステムによるメンテナンスコスト増加",
        "• オンプレミス環境でのリソース制約",
        "• モバイル対応していないため外出先での業務に制約"
    ]
    apply_body_style(current_box, current_state)
    
    # 右側: 課題パネル
    challenge_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.5),
        Inches(5.5), Inches(5),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.LIGHT_GRAY, line_width=Pt(0.75),
        shadow=True
    )
    
    # 課題パネルのヘッダー
    challenge_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(6.8), Inches(1.5),
        Inches(5.5), Inches(0.7),
        fill_color=ColorPalette.ACCENT2,
        text="解決すべき課題"
    )
    
    # 課題の内容
    challenge_box = slide.shapes.add_textbox(
        Inches(7.1), Inches(2.4),
        Inches(5), Inches(4)
    )
    
    challenges = [
        "• データの一元管理と業務プロセスの標準化",
        "• システム間連携の自動化による二重作業の排除",
        "• クラウド環境への移行によるコスト最適化",
        "• モバイル対応によるリモートワーク環境の整備",
        "• セキュリティ強化とコンプライアンス対応"
    ]
    apply_body_style(challenge_box, challenges)
    
    # 中央の接続要素 (矢印)
    connector_shape = add_shape(
        slide, MSO_SHAPE.RIGHT_ARROW,
        Inches(6), Inches(4),
        Inches(1), Inches(0.4),
        fill_color=ColorPalette.GRAY
    )
    
    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="3/10")

def create_proposal(prs):
    """提案内容のスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs, "solid", ColorPalette.SECONDARY)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2),
        fill_color=ColorPalette.PRIMARY
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(1), Inches(0.3),
        Inches(11), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "提案内容: クラウド統合管理システム"
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.LIGHT
    
    # 装飾的な図形
    decorative_shape = add_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(10), Inches(5.5),
        Inches(3), Inches(3),
        fill_color=ColorPalette.ACCENT3, transparency=0.9
    )
    
    # メインコンテンツパネル
    main_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(7), Inches(5.5),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.LIGHT_GRAY, line_width=Pt(0.75),
        shadow=True
    )
    
    # システム概要パネル
    system_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.3), Inches(1.5),
        Inches(4), Inches(2.5),
        fill_color=ColorPalette.ACCENT4, transparency=0.1,
        line_color=ColorPalette.ACCENT4, line_width=Pt(0.75)
    )
    
    # システム概要ヘッダー
    system_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(8.3), Inches(1.5),
        Inches(4), Inches(0.5),
        fill_color=ColorPalette.ACCENT4,
        text="システム概要"
    )
    
    # 主要機能パネル
    function_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.3), Inches(4.2),
        Inches(4), Inches(2.8),
        fill_color=ColorPalette.ACCENT1, transparency=0.1,
        line_color=ColorPalette.ACCENT1, line_width=Pt(0.75)
    )
    
    # 主要機能ヘッダー
    function_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(8.3), Inches(4.2),
        Inches(4), Inches(0.5),
        fill_color=ColorPalette.ACCENT1,
        text="主要機能"
    )
    
    # メインコンテンツ
    main_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2),
        Inches(6), Inches(4.8)
    )
    
    main_content = [
        "【クラウド統合管理システムの特徴】",
        "",
        "• すべての業務データを一元管理",
        "• どこからでもアクセス可能なクラウド基盤",
        "• 直感的なユーザーインターフェース",
        "• リアルタイムでのデータ同期と分析",
        "• 業務プロセスの自動化による効率化",
        "• AIによる予測分析と意思決定支援",
        "• 柔軟なスケーリングとカスタマイズ性"
    ]
    apply_body_style(main_box, main_content)
    
    # システム概要コンテンツ
    system_box = slide.shapes.add_textbox(
        Inches(8.5), Inches(2.1),
        Inches(3.6), Inches(1.8)
    )
    
    system_content = [
        "• クラウドベースの統合プラットフォーム",
        "• マイクロサービスアーキテクチャ",
        "• APIによる外部システム連携",
        "• レスポンシブデザイン対応",
        "• セキュアなデータストレージ"
    ]
    apply_body_style(system_box, system_content, BODY_SIZE, ColorPalette.DARK, Pt(8))
    
    # 主要機能コンテンツ
    function_box = slide.shapes.add_textbox(
        Inches(8.5), Inches(4.8),
        Inches(3.6), Inches(2)
    )
    
    function_content = [
        "• 顧客情報・案件管理の統合",
        "• リアルタイムダッシュボード",
        "• ワークフロー自動化",
        "• 権限管理とセキュリティ制御",
        "• モバイルアプリケーション"
    ]
    apply_body_style(function_box, function_content, BODY_SIZE, ColorPalette.DARK, Pt(8))
    
    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="4/10")

def create_schedule(prs):
    """導入スケジュールのスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs, "solid", ColorPalette.SECONDARY)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2),
        fill_color=ColorPalette.PRIMARY
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(1), Inches(0.3),
        Inches(11), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "導入スケジュール（6ヶ月計画）"
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.LIGHT
    
    # 装飾的な要素
    decorative_shape = add_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(11), Inches(1.5),
        Inches(2), Inches(2),
        fill_color=ColorPalette.ACCENT1, transparency=0.9
    )
    
    # タイムラインの作成
    # フェーズ1
    phase1_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(11), Inches(1.25),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.ACCENT1, line_width=Pt(1.5),
        shadow=True
    )
    
    # フェーズ1ヘッダー
    phase1_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(2.5), Inches(1.25),
        fill_color=ColorPalette.ACCENT1,
        text="フェーズ1\n要件定義・設計\n4月〜5月"
    )
    
    # フェーズ1内容
    phase1_content = slide.shapes.add_textbox(
        Inches(3.6), Inches(1.7),
        Inches(8), Inches(0.9)
    )
    
    phase1_text = [
        "• 業務要件の詳細ヒアリングと分析",
        "• システム設計とアーキテクチャ確定",
        "• データ移行計画の策定"
    ]
    apply_body_style(phase1_content, phase1_text, BODY_SIZE, ColorPalette.DARK, Pt(6))
    
    # フェーズ2
    phase2_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(2.9),
        Inches(11), Inches(1.25),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.ACCENT2, line_width=Pt(1.5),
        shadow=True
    )
    
    # フェーズ2ヘッダー
    phase2_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(2.9),
        Inches(2.5), Inches(1.25),
        fill_color=ColorPalette.ACCENT2,
        text="フェーズ2\n開発・構築\n5月〜7月"
    )
    
    # フェーズ2内容
    phase2_content = slide.shapes.add_textbox(
        Inches(3.6), Inches(3.1),
        Inches(8), Inches(0.9)
    )
    
    phase2_text = [
        "• システム基盤構築とコア機能の開発",
        "• 外部システム連携の実装",
        "• ユーザーインターフェース開発"
    ]
    apply_body_style(phase2_content, phase2_text, BODY_SIZE, ColorPalette.DARK, Pt(6))
    
    # フェーズ3
    phase3_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(4.3),
        Inches(11), Inches(1.25),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.ACCENT3, line_width=Pt(1.5),
        shadow=True
    )
    
    # フェーズ3ヘッダー
    phase3_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(4.3),
        Inches(2.5), Inches(1.25),
        fill_color=ColorPalette.ACCENT3,
        text="フェーズ3\nテスト・移行\n7月〜8月"
    )
    
    # フェーズ3内容
    phase3_content = slide.shapes.add_textbox(
        Inches(3.6), Inches(4.5),
        Inches(8), Inches(0.9)
    )
    
    phase3_text = [
        "• 単体・結合テストの実施",
        "• ユーザー受け入れテスト",
        "• データ移行とシステム切り替え準備"
    ]
    apply_body_style(phase3_content, phase3_text, BODY_SIZE, ColorPalette.DARK, Pt(6))
    
    # フェーズ4
    phase4_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(5.7),
        Inches(11), Inches(1.25),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.ACCENT4, line_width=Pt(1.5),
        shadow=True
    )
    
    # フェーズ4ヘッダー
    phase4_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(5.7),
        Inches(2.5), Inches(1.25),
        fill_color=ColorPalette.ACCENT4,
        text="フェーズ4\n本番稼働・安定化\n9月"
    )
    
    # フェーズ4内容
    phase4_content = slide.shapes.add_textbox(
        Inches(3.6), Inches(5.9),
        Inches(8), Inches(0.9)
    )
    
    phase4_text = [
        "• 段階的な本番リリース",
        "• ユーザートレーニングの実施",
        "• 運用体制の確立とサポート"
    ]
    apply_body_style(phase4_content, phase4_text, BODY_SIZE, ColorPalette.DARK, Pt(6))
    
    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="5/10")

def create_team_structure(prs):
    """実施体制のスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs, "solid", ColorPalette.SECONDARY)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2),
        fill_color=ColorPalette.PRIMARY
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(1), Inches(0.3),
        Inches(11), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "実施体制"
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.LIGHT
    
    # 装飾的な図形
    decorative_shape = add_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(0.5), Inches(6),
        Inches(1.5), Inches(1.5),
        fill_color=ColorPalette.ACCENT2, transparency=0.8
    )
    
    # プロジェクト推進体制パネル
    team_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(5), Inches(2.5),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.LIGHT_GRAY, line_width=Pt(0.75),
        shadow=True
    )
    
    # プロジェクト推進体制ヘッダー
    team_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(5), Inches(0.5),
        fill_color=ColorPalette.ACCENT1,
        text="プロジェクト推進体制"
    )
    
    # プロジェクト推進体制内容
    team_content = slide.shapes.add_textbox(
        Inches(1.3), Inches(2.1),
        Inches(4.5), Inches(1.8)
    )
    
    team_text = [
        "• プロジェクトスポンサー: 経営企画部長",
        "• プロジェクトマネージャー: IT部門 課長",
        "• テクニカルリード: システム開発チーム リーダー",
        "• 業務プロセス担当: 各部門代表者"
    ]
    apply_body_style(team_content, team_text, BODY_SIZE, ColorPalette.DARK, Pt(8))
    
    # 役割と責任パネル
    role_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.5), Inches(1.5),
        Inches(5.8), Inches(2.5),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.LIGHT_GRAY, line_width=Pt(0.75),
        shadow=True
    )
    
    # 役割と責任ヘッダー
    role_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(6.5), Inches(1.5),
        Inches(5.8), Inches(0.5),
        fill_color=ColorPalette.ACCENT2,
        text="役割と責任"
    )
    
    # 役割と責任内容
    role_content = slide.shapes.add_textbox(
        Inches(6.8), Inches(2.1),
        Inches(5.3), Inches(1.8)
    )
    
    role_text = [
        "• 要件定義・設計: 弊社コンサルタント + お客様業務担当者",
        "• システム開発: 弊社エンジニアチーム（5名）",
        "• テスト・品質保証: 弊社QAチーム + お客様検証担当者",
        "• 導入・トレーニング: 弊社導入支援チーム"
    ]
    apply_body_style(role_content, role_text, BODY_SIZE, ColorPalette.DARK, Pt(8))
    
    # コミュニケーション体制パネル
    comm_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(4.3),
        Inches(11.3), Inches(2.5),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.LIGHT_GRAY, line_width=Pt(0.75),
        shadow=True
    )
    
    # コミュニケーション体制ヘッダー
    comm_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(4.3),
        Inches(11.3), Inches(0.5),
        fill_color=ColorPalette.ACCENT3,
        text="コミュニケーション体制"
    )
    
    # コミュニケーション体制内容
    comm_content = slide.shapes.add_textbox(
        Inches(1.3), Inches(4.9),
        Inches(10.8), Inches(1.8)
    )
    
    comm_text = [
        "• 週次進捗会議（オンライン）: プロジェクトマネージャーが進捗、課題、リスクを報告",
        "• 月次ステアリングコミッティ（対面）: 経営層へ報告、重要決定事項の承認",
        "• 日次スクラムミーティング（開発チーム）: 15分の短時間で作業状況共有",
        "• 課題管理システムによるリアルタイム状況共有と透明性確保"
    ]
    apply_body_style(comm_content, comm_text, BODY_SIZE, ColorPalette.DARK, Pt(8))
    
    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="6/10")

def create_risk_management(prs):
    """リスク管理計画のスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs, "solid", ColorPalette.SECONDARY)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2),
        fill_color=ColorPalette.PRIMARY
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(1), Inches(0.3),
        Inches(11), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "リスク管理計画"
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.LIGHT
    
    # メインパネル
    main_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(11.3), Inches(5.3),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.LIGHT_GRAY, line_width=Pt(0.75),
        shadow=True
    )
    
    # リスク1パネル
    risk1_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.3), Inches(1.8),
        Inches(5.3), Inches(1.6),
        fill_color=ColorPalette.ACCENT1, transparency=0.9,
        line_color=ColorPalette.ACCENT1, line_width=Pt(0.75)
    )
    
    # リスク1内容
    risk1_title = slide.shapes.add_textbox(
        Inches(1.5), Inches(1.9),
        Inches(5), Inches(0.4)
    )
    risk1_tf = risk1_title.text_frame
    risk1_p = risk1_tf.paragraphs[0]
    risk1_p.text = "リスク1: 要件定義の不足・変更による開発遅延"
    risk1_run = risk1_p.runs[0]
    risk1_run.font.name = BODY_FONT
    risk1_run.font.size = BODY_SIZE
    risk1_run.font.bold = True
    risk1_run.font.color.rgb = ColorPalette.DARK
    
    risk1_content = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.3),
        Inches(5), Inches(1)
    )
    risk1_text = [
        "対策: ",
        "• アジャイル開発手法の採用",
        "• 定期的な要件レビュー会議の実施",
        "• 変更管理プロセスの厳格化"
    ]
    apply_body_style(risk1_content, risk1_text, BODY_SIZE, ColorPalette.DARK, Pt(6))
    
    # リスク2パネル
    risk2_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.8),
        Inches(5.3), Inches(1.6),
        fill_color=ColorPalette.ACCENT2, transparency=0.9,
        line_color=ColorPalette.ACCENT2, line_width=Pt(0.75)
    )
    
    # リスク2内容
    risk2_title = slide.shapes.add_textbox(
        Inches(7), Inches(1.9),
        Inches(5), Inches(0.4)
    )
    risk2_tf = risk2_title.text_frame
    risk2_p = risk2_tf.paragraphs[0]
    risk2_p.text = "リスク2: データ移行時のデータ欠損・不整合"
    risk2_run = risk2_p.runs[0]
    risk2_run.font.name = BODY_FONT
    risk2_run.font.size = BODY_SIZE
    risk2_run.font.bold = True
    risk2_run.font.color.rgb = ColorPalette.DARK
    
    risk2_content = slide.shapes.add_textbox(
        Inches(7), Inches(2.3),
        Inches(5), Inches(1)
    )
    risk2_text = [
        "対策: ",
        "• 事前データクレンジングの実施",
        "• 段階的移行アプローチの採用",
        "• 二重検証体制の構築"
    ]
    apply_body_style(risk2_content, risk2_text, BODY_SIZE, ColorPalette.DARK, Pt(6))
    
    # リスク3パネル
    risk3_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.3), Inches(3.6),
        Inches(5.3), Inches(1.6),
        fill_color=ColorPalette.ACCENT3, transparency=0.9,
        line_color=ColorPalette.ACCENT3, line_width=Pt(0.75)
    )
    
    # リスク3内容
    risk3_title = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.7),
        Inches(5), Inches(0.4)
    )
    risk3_tf = risk3_title.text_frame
    risk3_p = risk3_tf.paragraphs[0]
    risk3_p.text = "リスク3: ユーザー受け入れの低さ"
    risk3_run = risk3_p.runs[0]
    risk3_run.font.name = BODY_FONT
    risk3_run.font.size = BODY_SIZE
    risk3_run.font.bold = True
    risk3_run.font.color.rgb = ColorPalette.DARK
    
    risk3_content = slide.shapes.add_textbox(
        Inches(1.5), Inches(4.1),
        Inches(5), Inches(1)
    )
    risk3_text = [
        "対策: ",
        "• 早期からのユーザー参加促進",
        "• 充実したトレーニング計画の策定",
        "• ユーザーフィードバックの継続的収集"
    ]
    apply_body_style(risk3_content, risk3_text, BODY_SIZE, ColorPalette.DARK, Pt(6))
    
    # リスク4パネル
    risk4_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(3.6),
        Inches(5.3), Inches(1.6),
        fill_color=ColorPalette.ACCENT4, transparency=0.9,
        line_color=ColorPalette.ACCENT4, line_width=Pt(0.75)
    )
    
    # リスク4内容
    risk4_title = slide.shapes.add_textbox(
        Inches(7), Inches(3.7),
        Inches(5), Inches(0.4)
    )
    risk4_tf = risk4_title.text_frame
    risk4_p = risk4_tf.paragraphs[0]
    risk4_p.text = "リスク4: 既存システムとの連携不具合"
    risk4_run = risk4_p.runs[0]
    risk4_run.font.name = BODY_FONT
    risk4_run.font.size = BODY_SIZE
    risk4_run.font.bold = True
    risk4_run.font.color.rgb = ColorPalette.DARK
    
    risk4_content = slide.shapes.add_textbox(
        Inches(7), Inches(4.1),
        Inches(5), Inches(1)
    )
    risk4_text = [
        "対策: ",
        "• 詳細なインターフェース設計",
        "• 段階的な連携テストの実施",
        "• フォールバック機構の設計"
    ]
    apply_body_style(risk4_content, risk4_text, BODY_SIZE, ColorPalette.DARK, Pt(6))
    
    # リスク5パネル
    risk5_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4), Inches(5.3),
        Inches(5.3), Inches(1.3),
        fill_color=ColorPalette.PRIMARY, transparency=0.85,
        line_color=ColorPalette.PRIMARY, line_width=Pt(0.75)
    )
    
    # リスク5内容
    risk5_title = slide.shapes.add_textbox(
        Inches(4.2), Inches(5.4),
        Inches(5), Inches(0.4)
    )
    risk5_tf = risk5_title.text_frame
    risk5_p = risk5_tf.paragraphs[0]
    risk5_p.text = "リスク5: セキュリティインシデント"
    risk5_run = risk5_p.runs[0]
    risk5_run.font.name = BODY_FONT
    risk5_run.font.size = BODY_SIZE
    risk5_run.font.bold = True
    risk5_run.font.color.rgb = ColorPalette.LIGHT
    
    risk5_content = slide.shapes.add_textbox(
        Inches(4.2), Inches(5.8),
        Inches(5), Inches(0.7)
    )
    risk5_text = [
        "対策: セキュリティ設計レビュー、脆弱性診断、インシデント対応計画の策定"
    ]
    apply_body_style(risk5_content, risk5_text, BODY_SIZE, ColorPalette.LIGHT, Pt(6))
    
    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="7/10")

def create_budget(prs):
    """予算計画のスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs, "solid", ColorPalette.SECONDARY)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2),
        fill_color=ColorPalette.PRIMARY
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(1), Inches(0.3),
        Inches(11), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "予算計画"
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.LIGHT
    
    # 装飾的な要素
    decorative_shape = add_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(10.5), Inches(5.5),
        Inches(3), Inches(3),
        fill_color=ColorPalette.ACCENT3, transparency=0.9
    )
    
    # 初期導入コストパネル
    initial_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(6), Inches(2.8),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.LIGHT_GRAY, line_width=Pt(0.75),
        shadow=True
    )
    
    # 初期導入コストヘッダー
    initial_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(6), Inches(0.5),
        fill_color=ColorPalette.ACCENT1,
        text="初期導入コスト"
    )
    
    # 初期導入コスト内容
    initial_content = slide.shapes.add_textbox(
        Inches(1.3), Inches(2.1),
        Inches(5.4), Inches(2)
    )
    
    initial_text = [
        "• システム設計・開発費: 2,000万円",
        "• ハードウェア・クラウド環境構築: 500万円",
        "• データ移行・テスト: 600万円",
        "• トレーニング・導入支援: 400万円",
        "• 初期費用合計: 3,500万円"
    ]
    apply_body_style(initial_content, initial_text, BODY_SIZE, ColorPalette.DARK, Pt(8))
    
    # ランニングコストパネル
    running_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(4.5),
        Inches(6), Inches(2.3),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.LIGHT_GRAY, line_width=Pt(0.75),
        shadow=True
    )
    
    # ランニングコストヘッダー
    running_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(4.5),
        Inches(6), Inches(0.5),
        fill_color=ColorPalette.ACCENT2,
        text="ランニングコスト（年間）"
    )
    
    # ランニングコスト内容
    running_content = slide.shapes.add_textbox(
        Inches(1.3), Inches(5.1),
        Inches(5.4), Inches(1.6)
    )
    
    running_text = [
        "• クラウドインフラ利用料: 300万円",
        "• ライセンス費用: 200万円",
        "• 保守・サポート費: 300万円",
        "• 年間運用コスト合計: 800万円"
    ]
    apply_body_style(running_content, running_text, BODY_SIZE, ColorPalette.DARK, Pt(8))
    
    # ROIパネル
    roi_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.3), Inches(1.5),
        Inches(5), Inches(3.8),
        fill_color=ColorPalette.ACCENT4, transparency=0.1,
        line_color=ColorPalette.ACCENT4, line_width=Pt(0.75),
        shadow=True
    )
    
    # ROIヘッダー
    roi_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(7.3), Inches(1.5),
        Inches(5), Inches(0.5),
        fill_color=ColorPalette.ACCENT4,
        text="投資対効果（ROI）"
    )
    
    # ROI内容
    roi_content = slide.shapes.add_textbox(
        Inches(7.6), Inches(2.1),
        Inches(4.5), Inches(3)
    )
    
    roi_text = [
        "【コスト削減効果】",
        "• 業務効率化による人件費削減: 年間1,200万円",
        "• システム統合によるコスト削減: 年間800万円",
        "",
        "【定性的効果】",
        "• 意思決定スピードの向上",
        "• 顧客満足度の向上",
        "• データ活用による戦略的優位性",
        "",
        "【投資回収期間】",
        "• 約18ヶ月で初期投資を回収"
    ]
    apply_body_style(roi_content, roi_text, BODY_SIZE, ColorPalette.DARK, Pt(8))
    
    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="8/10")

def create_success_criteria(prs):
    """成功基準と評価方法のスライド"""
    slide_layout = prs.slide_layouts[1]  # タイトルと内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景設定
    add_background(slide, prs, "solid", ColorPalette.SECONDARY)
    
    # ヘッダーバー
    header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2),
        fill_color=ColorPalette.PRIMARY
    )
    
    # ヘッダータイトル
    header_title = slide.shapes.add_textbox(
        Inches(1), Inches(0.3),
        Inches(11), Inches(0.6)
    )
    header_tf = header_title.text_frame
    header_p = header_tf.paragraphs[0]
    header_p.text = "成功基準と評価方法"
    header_run = header_p.runs[0]
    header_run.font.name = TITLE_FONT
    header_run.font.size = HEADING_SIZE
    header_run.font.bold = True
    header_run.font.color.rgb = ColorPalette.LIGHT
    
    # 装飾的な要素
    decorative_shape = add_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(0.5), Inches(5.5),
        Inches(2), Inches(2),
        fill_color=ColorPalette.ACCENT1, transparency=0.9
    )
    
    # 成功基準メインパネル
    main_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5),
        Inches(11.3), Inches(2.5),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=ColorPalette.LIGHT_GRAY, line_width=Pt(0.75),
        shadow=True
    )
    
    # パフォーマンス指標パネル
    perf_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.3), Inches(1.8),
        Inches(5.3), Inches(2),
        fill_color=ColorPalette.ACCENT1, transparency=0.9,
        line_color=ColorPalette.ACCENT1, line_width=Pt(0.75)
    )
    
    # パフォーマンス指標ヘッダー
    perf_title = slide.shapes.add_textbox(
        Inches(1.5), Inches(1.9),
        Inches(5), Inches(0.4)
    )
    perf_tf = perf_title.text_frame
    perf_p = perf_tf.paragraphs[0]
    perf_p.text = "システムパフォーマンス指標"
    perf_run = perf_p.runs[0]
    perf_run.font.name = BODY_FONT
    perf_run.font.size = BODY_SIZE
    perf_run.font.bold = True
    perf_run.font.color.rgb = ColorPalette.DARK
    
    # パフォーマンス指標内容
    perf_content = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.3),
        Inches(5), Inches(1.4)
    )
    perf_text = [
        "• システム応答時間: 2秒以内（ピーク時）",
        "• システム可用性: 99.9%以上",
        "• 同時接続ユーザー: 最大300名をサポート",
        "• バックアップ復旧時間: 4時間以内"
    ]
    apply_body_style(perf_content, perf_text, BODY_SIZE, ColorPalette.DARK, Pt(6))
    
    # ビジネス効果パネル
    biz_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.8),
        Inches(5.3), Inches(2),
        fill_color=ColorPalette.ACCENT2, transparency=0.9,
        line_color=ColorPalette.ACCENT2, line_width=Pt(0.75)
    )
    
    # ビジネス効果ヘッダー
    biz_title = slide.shapes.add_textbox(
        Inches(7), Inches(1.9),
        Inches(5), Inches(0.4)
    )
    biz_tf = biz_title.text_frame
    biz_p = biz_tf.paragraphs[0]
    biz_p.text = "ビジネス効果指標"
    biz_run = biz_p.runs[0]
    biz_run.font.name = BODY_FONT
    biz_run.font.size = BODY_SIZE
    biz_run.font.bold = True
    biz_run.font.color.rgb = ColorPalette.DARK
    
    # ビジネス効果内容
    biz_content = slide.shapes.add_textbox(
        Inches(7), Inches(2.3),
        Inches(5), Inches(1.4)
    )
    biz_text = [
        "• 業務処理時間: 30%削減",
        "• 顧客対応時間: 50%短縮",
        "• データ入力エラー: 90%削減",
        "• ペーパーレス化: 紙使用量80%削減",
        "• ユーザー満足度: 80%以上"
    ]
    apply_body_style(biz_content, biz_text, BODY_SIZE, ColorPalette.DARK, Pt(6))
    
    # 評価方法パネル
    eval_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(4.2),
        Inches(11.3), Inches(2.5),
        fill_color=ColorPalette.ACCENT3, transparency=0.1,
        line_color=ColorPalette.ACCENT3, line_width=Pt(0.75),
        shadow=True
    )
    
    # 評価方法ヘッダー
    eval_header = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(4.2),
        Inches(11.3), Inches(0.5),
        fill_color=ColorPalette.ACCENT3,
        text="測定・評価方法"
    )
    
    # 評価方法内容
    eval_content = slide.shapes.add_textbox(
        Inches(1.3), Inches(4.8),
        Inches(10.8), Inches(1.8)
    )
    
    eval_text = [
        "【定期評価】",
        "• 四半期ごとのパフォーマンス測定レポート作成",
        "• 月次ユーザー満足度調査の実施",
        "• 業務効率化指標の定期測定と分析",
        "",
        "【継続的モニタリング】",
        "• リアルタイムダッシュボードによるシステム状態監視",
        "• インシデント発生率とレスポンス時間の追跡",
        "• ユーザーフィードバックの継続的収集と分析"
    ]
    apply_body_style(eval_content, eval_text, BODY_SIZE, ColorPalette.DARK, Pt(8))
    
    # フッター追加
    add_footer(slide, prs, include_page_number=True, page_num="9/10")

def create_conclusion(prs):
    """まとめと次のステップのスライド"""
    slide_layout = prs.slide_layouts[0]  # タイトルスライド
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景全体にグラデーション
    background = add_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height,
        fill_color=ColorPalette.PRIMARY, gradient_to=ColorPalette.DARK
    )
    
    # 装飾的な要素
    decorative_shape1 = add_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(1), Inches(1),
        Inches(2), Inches(2),
        fill_color=ColorPalette.ACCENT1, transparency=0.85
    )
    
    decorative_shape2 = add_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(10), Inches(5),
        Inches(3), Inches(3),
        fill_color=ColorPalette.ACCENT2, transparency=0.85
    )
    
    # まとめタイトル
    title_box = slide.shapes.add_textbox(
        Inches(2.5), Inches(1.5),
        Inches(8), Inches(1)
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "まとめと次のステップ"
    title_p.alignment = PP_ALIGN.CENTER
    title_run = title_p.runs[0]
    title_run.font.name = TITLE_FONT
    title_run.font.size = HEADING_SIZE
    title_run.font.bold = True
    title_run.font.color.rgb = ColorPalette.LIGHT
    
    # まとめパネル
    summary_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.5), Inches(2.5),
        Inches(8), Inches(1.8),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=None, shadow=True
    )
    
    # まとめ内容
    summary_content = slide.shapes.add_textbox(
        Inches(2.8), Inches(2.7),
        Inches(7.4), Inches(1.4)
    )
    
    summary_text = [
        "【提案のまとめ】",
        "• クラウドベースの統合管理システム導入により業務効率を30%向上",
        "• 6ヶ月間の段階的な導入計画で業務への影響を最小化",
        "• 初期投資3,500万円、年間運用コスト800万円、18ヶ月でROI達成"
    ]
    apply_body_style(summary_content, summary_text, BODY_SIZE, ColorPalette.DARK, Pt(8))
    
    # 次のステップパネル
    next_panel = add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.5), Inches(4.5),
        Inches(8), Inches(1.8),
        fill_color=ColorPalette.LIGHT, transparency=0,
        line_color=None, shadow=True
    )
    
    # 次のステップ内容
    next_content = slide.shapes.add_textbox(
        Inches(2.8), Inches(4.7),
        Inches(7.4), Inches(1.4)
    )
    
    next_text = [
        "【次のステップ】",
        "• 提案内容の最終確認と承認（1週間以内）",
        "• キックオフミーティングの開催（承認後2週間以内）",
        "• 詳細要件定義の開始（4月第1週）"
    ]
    apply_body_style(next_content, next_text, BODY_SIZE, ColorPalette.DARK, Pt(8))
    
    # 連絡先
    contact_box = slide.shapes.add_textbox(
        Inches(2.5), Inches(6.5),
        Inches(8), Inches(0.8)
    )
    contact_tf = contact_box.text_frame
    contact_p = contact_tf.paragraphs[0]
    contact_p.text = "プロジェクト担当: 山田太郎 | yamada.taro@example.com | 03-1234-5678"
    contact_p.alignment = PP_ALIGN.CENTER
    contact_run = contact_p.runs[0]
    contact_run.font.name = BODY_FONT
    contact_run.font.size = BODY_SIZE
    contact_run.font.color.rgb = ColorPalette.LIGHT

if __name__ == "__main__":
    create_presentation()